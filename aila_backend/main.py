from dotenv import load_dotenv
load_dotenv()
import os

from fastapi import FastAPI, Depends, HTTPException, UploadFile, File, Form, BackgroundTasks, Body
from fastapi import WebSocket, WebSocketDisconnect, Path
from fastapi.middleware.cors import CORSMiddleware

from sqlalchemy.orm import Session
from sqlalchemy import Column, Integer, String, DateTime, ForeignKey, Text, JSON, Boolean, func, Index, PrimaryKeyConstraint
from sqlalchemy.sql import func
from sqlalchemy import Boolean
from sqlalchemy import text as sql_text
from sqlalchemy import or_

from pydantic import BaseModel, Field
from llama_index.llms.gemini import Gemini
from pptx import Presentation
from typing import List, Optional, Dict, Set
from collections import deque, defaultdict, Counter
from datetime import datetime

import asyncio
import fitz  # PyMuPDF
import hashlib
import traceback
import random
import uuid
import shutil
import json
import re
import logging
import google.generativeai as genai

from aila_backend.database import SessionLocal, Base, engine
from aila_backend.models import (
    User,
    Course,
    Enrollment,
    LectureUpload,
    LectureProcessing,
    Segment,
    KnowledgeGraph,
    Quiz,
    MCQ,
    QuizSettings,
    QuizAttempt,
    MCQResponse,
)

UPLOAD_DIR = "db/uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)
genai.configure(api_key=os.environ["GOOGLE_API_KEY"])

Base.metadata.create_all(bind=engine)

app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "http://127.0.0.1:3000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class MethodFilter(logging.Filter):
    def filter(self, record):
        if record.args and len(record.args) >= 4:
            method = record.args[1]
            status = str(record.args[3])
            if method in ("POST", "PUT", "DELETE", "PATCH"):
                return True
            if status.startswith(("4", "5")):
                return True
            return False
        return False

uvicorn_access_logger = logging.getLogger("uvicorn.access")
uvicorn_access_logger.addFilter(MethodFilter())


# === Pydantic Models (for API request/response validation) ===

class MCQReqModel(BaseModel):
    course_id: str
    week: int
    concept_id: str
    summary: str = ""
    contents: str = ""


class MCQPayload(BaseModel):
    segment_id: str
    content: str


class MCQConceptModel(BaseModel):
    course_id: str
    week: int
    concept_id: str
    summary: str = ""
    contents: str = ""


class QuizCreate(BaseModel):
    name: str = "New Quiz"
    course_id: str
    week: int
    instructor_id: str = "unknown"
    concept_ids: List[str]


class QuizSettingsIn(BaseModel):
    week: int
    min_difficulty: str = "Easy"
    max_difficulty: str = "Hard"
    min_bloom_level: str = "Remember" 
    max_bloom_level: str = "Create"    
    max_questions: int = 10
    allowed_retries: int = 3
    feedback_style: str = "Immediate"
    include_spaced: bool = False

class QuizSettingsOut(BaseModel):
    id: int
    quiz_id: str
    quizid: str
    week: int
    mindifficulty: str
    maxdifficulty: str
    min_bloom_level: str  
    max_bloom_level: str  
    maxquestions: int
    allowedretries: int
    feedbackstyle: str
    includespaced: bool

    class Config:
        from_attributes = True


class QuizStartResponse(BaseModel):
    attempt_id: str
    questions: List[dict]
    settings: dict
    retries_left: int


class MCQUpdate(BaseModel):
    question: str
    options: List[str]
    answer: str
    difficulty: Optional[str] = "Medium"
    bloom_level: Optional[str] = "Remember"


class GenerateQuizMCQsRequest(BaseModel):
    action: str = "generate_from_concepts"


class TitleGenRequest(BaseModel):
    concepts: List[str]


class LockPreviewPayload(BaseModel):
    quiz_id: str
    items: List[dict]


# === WebSocket Manager ===

class ConnectionManager:
    def __init__(self):
        self.active_connections: Dict[str, Set[WebSocket]] = {}

    async def connect(self, websocket: WebSocket, quiz_id: str):
        await websocket.accept()
        if quiz_id not in self.active_connections:
            self.active_connections[quiz_id] = set()
        self.active_connections[quiz_id].add(websocket)
        print(f"✅ Instructor connected to quiz {quiz_id}")

    def disconnect(self, websocket: WebSocket, quiz_id: str):
        if quiz_id in self.active_connections:
            self.active_connections[quiz_id].discard(websocket)
            if not self.active_connections[quiz_id]:
                del self.active_connections[quiz_id]
        print(f"❌ Instructor disconnected from quiz {quiz_id}")

    async def broadcast(self, quiz_id: str, message: dict):
        if quiz_id in self.active_connections:
            dead_connections = set()
            for connection in self.active_connections[quiz_id]:
                try:
                    await connection.send_json(message)
                except Exception as e:
                    print(f"Failed to send to connection: {e}")
                    dead_connections.add(connection)
            
            for conn in dead_connections:
                self.active_connections[quiz_id].discard(conn)

manager = ConnectionManager()


def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()


def normalize_id(text):
    """Make sure every concept has the EXACT same ID everywhere"""
    return re.sub(r'\W+', '_', text.strip().lower()).strip('_')

# --- HELPER: ROBUST JSON PARSER ---
def repair_json(json_str):
    """
    Attempts to fix common JSON errors from LLMs.
    """
    try:
        return json.loads(json_str)
    except:
        pass
        
    # Remove markdown code blocks
    clean = json_str.replace("```json", "").replace("```", "").strip()
    
    # Try finding the outer braces
    start = clean.find("{")
    end = clean.rfind("}") + 1
    if start != -1 and end != 0:
        clean = clean[start:end]
        
    try:
        return json.loads(clean)
    except:
        # Last resort: Try to fix trailing commas
        try:
            clean = re.sub(r",\s*}", "}", clean)
            clean = re.sub(r",\s*]", "]", clean)
            return json.loads(clean)
        except:
            print(f"[JSON ERROR] Failed to parse: {clean[:100]}...")
            return {}

# --- HELPER: STEP 1 - STRUCTURE IDENTIFICATION ---
def identify_structure(llm, full_text, file_name):
    """
    First Pass: Identify the high-level outline AND anchor each sub-topic
    to the exact slide range where it appears.
    """
    print(f"🧠 [PASS 1] Identifying structure for {file_name}...")

    prompt = f"""
    Analyze lecture slides "{file_name}".

    GOAL: Identify the main topic and every distinct top-level sub-topic.

    CRITICAL RULES:
    1. Sub-topics must be the CORE NAMED CONCEPTS, not examples or details.
       - Lecture "4 Principles of OOP" → sub_topics: ["Encapsulation", "Inheritance", "Polymorphism", "Abstraction"]
       - Lecture "Stacks and Queues" → sub_topics: ["Stacks", "Queues"]
       - Lecture "Sorting Algorithms" → sub_topics: ["Bubble Sort", "Merge Sort", "Quick Sort"]
    2. NEVER use examples, code snippets, or slides as sub-topics.
    3. NEVER list "Introduction", "Summary", "Overview", "Review" as sub-topics.
    4. For each sub-topic, find the slide numbers where it is PRIMARILY discussed.

    Return JSON (no extra text, no markdown):
    {{
        "main_topic": "The single overarching subject (e.g. Object-Oriented Programming)",
        "sub_topics": [
            {{"name": "Encapsulation", "slide_nums": [3, 4, 5]}},
            {{"name": "Inheritance",   "slide_nums": [6, 7, 8]}}
        ]
    }}

    Context:
    {full_text[:16000]}
    """

    resp = llm.complete(prompt)
    raw = repair_json(str(resp))

    # Normalise: support both old format (list of strings) and new format (list of dicts)
    sub_topics_raw = raw.get("sub_topics", [])
    sub_topics_normalised = []
    for st in sub_topics_raw:
        if isinstance(st, str):
            sub_topics_normalised.append({"name": st, "slide_nums": []})
        elif isinstance(st, dict) and "name" in st:
            sub_topics_normalised.append(st)
    raw["sub_topics"] = sub_topics_normalised
    return raw

# --- HELPER: STEP 2a - SINGLE SUB-TOPIC EXTRACTION ---
def extract_concepts_for_subtopic(llm, main_topic, sub_topic_name, sub_topic_id, text_slice):
    """
    Extract child concepts for ONE sub-topic.
    sub_topic_id is the normalised snake_case ID we've already created as a node.
    The LLM only needs to return CHILDREN of that node.
    """
    prompt = f"""
    You are an expert Computer Science Curriculum Designer.

    CONTEXT: We are building a concept map for "{main_topic}".
    The sub-topic "{sub_topic_name}" is already a node.
    Your job: find its KEY CHILD CONCEPTS only.

    RULES:
    1. Return 2-4 child concepts MAX. These must be direct sub-concepts of "{sub_topic_name}".
    2. Children must be CONCEPTUAL (definitions, mechanisms, properties) — NOT examples, NOT code.
       - GOOD child of "Encapsulation": "Access Modifiers", "Data Hiding", "Getter/Setter Methods"
       - BAD child of "Encapsulation": "BankAccount class", "private int balance", "Python example"
    3. MERGE related details into one node:
       - BAD: "public", "private", "protected" (3 nodes) → GOOD: "Access Modifiers" (1 node)
    4. BANNED: "Introduction", "Summary", "Example", "Overview", "Review", code literals.
    5. Edge relation must be one of:
       is_a, has_part, uses, implements, requires, produces, example_of, contrasts_with, precedes

    Return ONLY this JSON structure (no markdown, no extra text):
    {{
        "nodes": [
            {{
                "id": "snake_case_unique_id",
                "label": "Human Readable Label",
                "type": "concept|structure|algorithm|detail",
                "summary": "One clear sentence explaining this concept.",
                "slide_nums": [4, 5]
            }}
        ],
        "edges": [
            {{ "source": "{sub_topic_id}", "target": "snake_case_unique_id", "relation": "has_part" }}
        ]
    }}

    Lecture text (focus on "{sub_topic_name}"):
    {text_slice[:10000]}
    """
    resp = llm.complete(prompt)
    result = repair_json(str(resp))
    nodes = result.get('nodes', [])
    edges = result.get('edges', [])
    print(f"  ✓ [{sub_topic_name}] → {len(nodes)} child nodes, {len(edges)} edges")
    return result


def _get_slide_anchored_slice(full_text, slide_nums, fallback_keyword, window=10000):
    """
    Build the best possible text slice for a sub-topic.
    Priority: use the slide numbers from identify_structure.
    Fallback: keyword search.
    """
    if slide_nums:
        # Extract the blocks for those specific slide numbers from full_text
        # Slide blocks are formatted as: "--- Title (Slide N) ---\n..."
        import re
        blocks = re.split(r'(?=--- .+ \(Slide \d+\) ---)', full_text)
        matched = []
        for block in blocks:
            m = re.match(r'--- .+ \(Slide (\d+)\) ---', block)
            if m and int(m.group(1)) in slide_nums:
                matched.append(block)
        if matched:
            result = "\n\n".join(matched)
            return result[:window]

    # Fallback: find keyword in text
    idx = full_text.lower().find(fallback_keyword.lower())
    if idx == -1:
        return full_text[:window]
    start = max(0, idx - 300)
    return full_text[start: start + window]


# --- HELPER: STEP 2 - CONCEPT EXTRACTION (parallel per sub-topic) ---
def extract_concepts(llm, structure, full_text):
    main_topic  = structure.get("main_topic", "Lecture")
    sub_topics  = structure.get("sub_topics", [])  # list of {name, slide_nums}

    print(f"🧠 [PASS 2] Extracting concepts for {len(sub_topics)} sub-topics in parallel...")

    if not sub_topics:
        sub_topics = [{"name": main_topic, "slide_nums": []}]

    # Pre-create a node for every sub-topic so they always exist in the graph
    # and the LLM only needs to return their children.
    sub_topic_nodes = []
    sub_topic_id_map = {}  # name -> id
    for st in sub_topics:
        st_name = st["name"]
        st_id   = normalize_id(st_name)  # e.g. "Encapsulation" → "encapsulation"
        sub_topic_id_map[st_name] = st_id
        sub_topic_nodes.append({
            "id":        st_id,
            "label":     st_name,
            "type":      "structure",
            "summary":   f"{st_name} is a core concept within {main_topic}.",
            "slide_nums": st.get("slide_nums", []),
            "isRoot":    False
        })

    # Run child-concept extractions in parallel
    from concurrent.futures import ThreadPoolExecutor, as_completed
    results_by_topic = {}
    with ThreadPoolExecutor(max_workers=min(len(sub_topics), 5)) as pool:
        future_to_name = {
            pool.submit(
                extract_concepts_for_subtopic,
                llm,
                main_topic,
                st["name"],
                sub_topic_id_map[st["name"]],
                _get_slide_anchored_slice(full_text, st.get("slide_nums", []), st["name"])
            ): st["name"]
            for st in sub_topics
        }
        for future in as_completed(future_to_name):
            st_name = future_to_name[future]
            try:
                results_by_topic[st_name] = future.result()
            except Exception as e:
                print(f"  ✗ [{st_name}] extraction failed: {e}")
                results_by_topic[st_name] = {"nodes": [], "edges": []}

    # Build final graph
    merged_nodes = []
    merged_edges = []
    seen_node_ids = set()

    # 1. Main topic root node
    main_id = normalize_id(main_topic)
    merged_nodes.append({
        "id":      main_id,
        "label":   main_topic,
        "type":    "root",
        "summary": f"Overview of {main_topic}.",
        "slide_nums": [],
        "isRoot":  True
    })
    seen_node_ids.add(main_id)

    # 2. Sub-topic nodes (always present, at level 1)
    for st_node in sub_topic_nodes:
        if st_node["id"] not in seen_node_ids:
            merged_nodes.append(st_node)
            seen_node_ids.add(st_node["id"])
        # Wire: main_topic → sub_topic
        merged_edges.append({
            "source":   main_id,
            "target":   st_node["id"],
            "relation": "has_part"
        })

    # 3. Child concept nodes from LLM
    for st_name, data in results_by_topic.items():
        for node in data.get("nodes", []):
            nid = node.get("id", "")
            if not nid:
                continue
            # Avoid overwriting a sub-topic node with a child node of the same ID
            if nid in seen_node_ids:
                continue
            # Reject nodes that look like the sub-topic itself (LLM sometimes echoes it)
            if _similar_labels(node.get("label", ""), st_name):
                continue
            seen_node_ids.add(nid)
            merged_nodes.append(node)
        for edge in data.get("edges", []):
            src = edge.get("source", "")
            tgt = edge.get("target", "")
            if src and tgt and src != tgt:
                merged_edges.append(edge)

    print(f"📊 [PASS 2 DONE] {len(merged_nodes)} nodes ({len(sub_topics)} sub-topics + {len(merged_nodes)-len(sub_topics)-1} children), {len(merged_edges)} edges")
    return {"nodes": merged_nodes, "edges": merged_edges}


# --- MAIN PIPELINE FUNCTION ---
def process_lecture_and_kg(filepath, upload_id, course_id, week, file_name, processing_id):
    db = SessionLocal()
    llm = Gemini(model="models/gemini-2.5-flash") 

    try:
        # 0. Init Status
        db.query(LectureProcessing).filter(LectureProcessing.id == processing_id).update({
            "status": "processing", "progress": 5, "error_message": None
        })
        db.commit()

        # ---------- 1. Extract Slides (Robust Text Extraction) ----------
        segments_data = []
        ext = os.path.splitext(filepath)[-1].lower()
        
        if ext == ".pdf":
            doc = fitz.open(filepath)
            for i, page in enumerate(doc):
                text = page.get_text(sort=True)
                if not text.strip():
                    continue
                # Extract title: largest-font span on the page
                page_title = f"Slide {i + 1}"
                try:
                    blocks = page.get_text("dict")["blocks"]
                    spans = [
                        s for b in blocks
                        for l in b.get("lines", [])
                        for s in l.get("spans", [])
                        if s.get("text", "").strip()
                    ]
                    if spans:
                        biggest = max(spans, key=lambda s: s["size"])
                        candidate = biggest["text"].strip()
                        if len(candidate) > 3:  # ignore single chars / page numbers
                            page_title = candidate
                except Exception:
                    pass
                segments_data.append({"slide_num": i + 1, "text": text, "title": page_title})
        elif ext == ".pptx":
            prs = Presentation(filepath)
            for i, slide in enumerate(prs.slides):
                slide_title = f"Slide {i + 1}"
                lines = []
                for shape in slide.shapes:
                    # Extract title from the first non-empty text frame
                    if hasattr(shape, "text") and shape.text.strip():
                        if slide_title == f"Slide {i + 1}":  # not yet set
                            slide_title = shape.text.strip().splitlines()[0][:120]
                        lines.append(shape.text)
                    # Also extract table cells
                    if shape.has_table:
                        for row_idx in range(shape.table.rows.count):
                            row_cells = [
                                shape.table.cell(row_idx, c).text.strip()
                                for c in range(shape.table.columns.count)
                            ]
                            row_text = "\t".join(cell for cell in row_cells if cell)
                            if row_text:
                                lines.append(row_text)
                content = "\n".join(lines)
                if content.strip():
                    segments_data.append({"slide_num": i + 1, "text": content, "title": slide_title})

        if not segments_data:
            raise ValueError("No readable slides found")

        # Save Segments
        print(f"[PROCESS] Saving {len(segments_data)} segments...")
        for seg in segments_data:
            new_seg = Segment(
                id=str(uuid.uuid4()),
                upload_id=upload_id,
                course_id=course_id,
                week=week,
                segment_index=seg["slide_num"],
                title=seg.get("title", f"Slide {seg['slide_num']}"),  # real title now
                content=seg["text"],
                keywords="", 
                summary=""
            )
            db.add(new_seg)
        db.commit()

        # Prepare Text — slide-boundary-aware truncation
        # Each slide is formatted with its real title so the LLM can see headings.
        def smart_truncate(segments, char_limit=28000):
            """Pack whole slides up to char_limit. Never cuts mid-slide."""
            budget = char_limit
            chunks = []
            for s in segments:
                slide_title = s.get("title", f"Slide {s['slide_num']}")
                block = f"--- {slide_title} (Slide {s['slide_num']}) ---\n{s['text']}"
                if len(block) > budget:
                    # If we have nothing yet, include a truncated version of the first slide
                    # rather than sending an empty string
                    if not chunks:
                        chunks.append(block[:budget])
                    break
                chunks.append(block)
                budget -= len(block)
            return "\n\n".join(chunks)

        full_text = smart_truncate(segments_data)
        total_chars = sum(len(s["text"]) for s in segments_data)
        if total_chars > 28000:
            covered = len(segments_data[:len(full_text.split("--- ")) - 1])
            print(f"⚠️ [WARN] Lecture {total_chars} chars. Fitting {len(full_text)} chars into LLM window (slide-boundary safe).")
        print(f"📄 [TEXT] {len(full_text)} chars across {len(segments_data)} slides")

        # ---------- 2. TWO-PASS GENERATION ----------
        
        # Pass 1: Identify Structure
        structure = identify_structure(llm, full_text, file_name)
        if not structure:
             # Fallback structure if LLM fails
             structure = {"main_topic": file_name, "sub_topics": []}
             
        # Pass 2: Extract Concepts based on Structure
        graph_data = extract_concepts(llm, structure, full_text)
        
        main_topic = structure.get("main_topic", file_name)
        # sub_topics is now list of {name, slide_nums} — extract names for logging
        st_names = [st["name"] if isinstance(st, dict) else st for st in structure.get("sub_topics", [])]
        print(f"[STRUCTURE] main='{main_topic}' sub_topics={st_names}")
        raw_nodes = graph_data.get("nodes", [])
        raw_edges = graph_data.get("edges", [])

        # ---------- 3. PROCESS & CLEAN ----------
        
        # Filter Junk
        clean_nodes = []
        clean_ids = set()
        for n in raw_nodes:
            lbl = n.get('label', '').lower()
            if any(x in lbl for x in ["init", "self.", "print(", "error", "summary", "intro"]):
                continue
            clean_nodes.append(n)
            clean_ids.add(n['id'])
            
        clean_edges = [e for e in raw_edges if e['source'] in clean_ids and e['target'] in clean_ids]
        
        print(f"📊 [GRAPH] Main: {main_topic} | Nodes: {len(clean_nodes)}")

        # Compute Hierarchy — use normalized ID as root (matches what extract_concepts creates)
        main_topic_id = normalize_id(main_topic)
        final_nodes, final_edges = compute_levels(clean_nodes, clean_edges, explicit_root=main_topic_id)

        # Restore Content Context
        slide_text_map = {s["slide_num"]: s["text"] for s in segments_data}
        for node in final_nodes:
            s_nums = node.get("slide_nums", [])
            gathered_text = []
            if isinstance(s_nums, (int, str)): s_nums = [int(s_nums)] if isinstance(s_nums, int) or (isinstance(s_nums, str) and s_nums.isdigit()) else []
            
            for num in s_nums:
                if num in slide_text_map:
                    gathered_text.append(slide_text_map[num])
            node["contents"] = "\n...\n".join(gathered_text)[:1500]

        # ---------- 4. SAVE & MERGE ----------
        
        # A. Save FILE Graph
        node_json = json.dumps(final_nodes)
        edge_json = json.dumps(final_edges)
        
        db.execute(
            sql_text("""
                INSERT INTO knowledge_graph (id, course_id, week, node_data, edge_data, graph_type, source_file) 
                VALUES (:id, :c, :w, :n, :e, 'file', :fname)
            """),
            {
                "id": str(uuid.uuid4()), "c": course_id, "w": week, 
                "n": node_json, "e": edge_json, "fname": file_name
            }
        )
        db.commit()

        # B. Merge
        print(f"🔄 [MERGE] Merging into Week {week} Master Graph...")
        
        all_files = db.execute(
            sql_text("SELECT node_data, edge_data FROM knowledge_graph WHERE course_id=:c AND week=:w AND graph_type='file'"),
            {"c": course_id, "w": week}
        ).fetchall()
        
        graphs_to_merge = []
        for row in all_files:
            try:
                graphs_to_merge.append({
                    "nodes": json.loads(row[0]),
                    "edges": json.loads(row[1])
                })
            except:
                continue
        
        # Use our Robust Merge
        master_data = merge_graphs(graphs_to_merge, week)
        
        # Update/Create Master
        master_node_json = json.dumps(master_data['nodes'])
        master_edge_json = json.dumps(master_data['edges'])
        
        existing_master = db.execute(
            sql_text("SELECT id FROM knowledge_graph WHERE course_id=:c AND week=:w AND graph_type='master'"),
            {"c": course_id, "w": week}
        ).fetchone()

        if existing_master:
            db.execute(
                sql_text("UPDATE knowledge_graph SET node_data=:n, edge_data=:e WHERE id=:id"),
                {"n": master_node_json, "e": master_edge_json, "id": existing_master[0]}
            )
        else:
            db.execute(
                sql_text("INSERT INTO knowledge_graph (id, course_id, week, node_data, edge_data, graph_type) VALUES (:id, :c, :w, :n, :e, 'master')"),
                {"id": str(uuid.uuid4()), "c": course_id, "w": week, "n": master_node_json, "e": master_edge_json}
            )
        db.commit()

        # Finish
        db.query(LectureProcessing).filter(LectureProcessing.id == processing_id).update({
            "status": "done", "progress": 100
        })
        db.commit()
        print(f"✅ [COMPLETE] Saved & Merged.")

    except Exception as e:
        print(f"❌ [ERROR] {str(e)}")
        traceback.print_exc()
        db.query(LectureProcessing).filter(LectureProcessing.id == processing_id).update({
            "status": "error", "progress": 0, "error_message": str(e)[:500]
        })
        db.commit()
    finally:
        db.close()




def compute_levels(nodes, edges, explicit_root=None):
    # 1. Normalize Node IDs (strip spaces, lower case check)
    # Create a map of "clean_label" -> "real_id"
    label_to_id = {}
    for n in nodes:
        label_to_id[n['id']] = n['id']
        if 'label' in n:
            label_to_id[n['label']] = n['id']
            # Add normalized versions for fuzzy matching
            label_to_id[n['label'].lower().strip()] = n['id']
            label_to_id[n['id'].lower().strip()] = n['id']

    # 2. Repair Edges (Fix ID mismatches)
    valid_edges = []
    for e in edges:
        src = e.get('source')
        tgt = e.get('target')
        
        # Try to find real IDs
        real_src = label_to_id.get(src) or label_to_id.get(src.lower().strip())
        real_tgt = label_to_id.get(tgt) or label_to_id.get(tgt.lower().strip())
        
        if real_src and real_tgt:
            valid_edges.append({
                "source": real_src, 
                "target": real_tgt, 
                "relation": e.get("relation", "related")
            })
    
    # Update edges list
    edges = valid_edges

    # 3. Ensure Explicit Root
    node_map = {n["id"]: n for n in nodes}
    if explicit_root:
        if explicit_root not in node_map:
            # Maybe the root ID changed during fuzzy match?
            # Try to find a node that looks like the root
            found = False
            for nid in node_map:
                if nid.lower() == explicit_root.lower():
                    explicit_root = nid
                    found = True
                    break
            
            if not found:
                # Create it if truly missing
                root_node = {"id": explicit_root, "label": explicit_root, "isRoot": True}
                nodes.append(root_node)
                node_map[explicit_root] = root_node
        
        node_map[explicit_root]["isRoot"] = True

    # 4. BFS (Standard)
    adj = defaultdict(list)
    for e in edges:
        adj[e["source"]].append(e["target"])

    levels = {node_id: -1 for node_id in node_map}
    queue = deque()
    
    if explicit_root and explicit_root in node_map:
        queue.append((explicit_root, 0))
    else:
        # Fallback for no explicit root
        incoming = defaultdict(int)
        for e in edges: incoming[e["target"]] += 1
        topo_roots = [n["id"] for n in nodes if incoming[n["id"]] == 0]
        for r in topo_roots: queue.append((r, 0))

    visited = set()
    if explicit_root: visited.add(explicit_root)

    while queue:
        curr, depth = queue.popleft()
        levels[curr] = depth
        for child in adj[curr]:
            if child not in visited:
                visited.add(child)
                queue.append((child, depth + 1))

    # 5. Rescue Orphans — try semantic parent first, fall back to root
    # Build list of already-reachable nodes for similarity comparison
    reachable_nodes = [n for n in nodes if levels.get(n["id"], -1) != -1]

    new_edges = []
    if explicit_root:
        for node in nodes:
            nid = node["id"]
            if nid == explicit_root or levels.get(nid, -1) != -1:
                continue  # already placed

            orphan_label = node.get("label", nid)

            # Try to find the best semantic parent among reachable nodes
            best_parent_id = explicit_root  # default fallback
            best_score = 0.0
            for candidate in reachable_nodes:
                if candidate["id"] == explicit_root:
                    continue
                candidate_label = candidate.get("label", candidate["id"])
                from difflib import SequenceMatcher
                score = SequenceMatcher(
                    None,
                    _normalize_label(orphan_label),
                    _normalize_label(candidate_label)
                ).ratio()
                if score > best_score:
                    best_score = score
                    best_parent_id = candidate["id"]

            # Only use the semantic parent if it's a reasonable match
            # (threshold 0.4 — loose, just avoids completely unrelated attachment)
            if best_score < 0.4:
                best_parent_id = explicit_root

            parent_level = levels.get(best_parent_id, 0)
            new_edges.append({
                "source": best_parent_id,
                "target": nid,
                "relation": "related"
            })
            levels[nid] = parent_level + 1
            node["isRoot"] = False
            # Add to reachable so subsequent orphans can attach to this one
            reachable_nodes.append(node)

    # 6. Finalize
    for node in nodes:
        node["level"] = levels.get(node["id"], 1)
        node["isRoot"] = (node["id"] == explicit_root)
        if node["isRoot"]: node["type"] = "root"
        elif "type" not in node: node["type"] = "concept"

    final_edges = edges + new_edges
    nodes.sort(key=lambda x: x["level"])
    
    return nodes, final_edges


    
def _normalize_label(label: str) -> str:
    """Lowercase, strip punctuation, collapse whitespace for similarity comparison."""
    import re
    return re.sub(r'[^a-z0-9 ]', '', label.lower()).strip()


def _similar_labels(a: str, b: str, threshold: float = 0.85) -> bool:
    """True if two concept labels are similar enough to be the same concept."""
    from difflib import SequenceMatcher
    na, nb = _normalize_label(a), _normalize_label(b)
    if na == nb:
        return True
    # Short labels need exact match; longer ones use similarity ratio
    if min(len(na), len(nb)) < 4:
        return na == nb
    return SequenceMatcher(None, na, nb).ratio() >= threshold


def merge_graphs(graphs_list, week_number):
    merged_nodes = {}     # id -> node dict
    id_remap = {}         # old_id -> canonical_id (for duplicate resolution)
    merged_edges = []
    seen_edges = set()

    # 1. Create Week Root (Level 0)
    week_root_id = f"Week {week_number} Overview"
    week_root = {
        "id": week_root_id,
        "label": f"Week {week_number} Concepts",
        "isRoot": True,
        "level": 0,
        "type": "root",
        "summary": "Overview of all topics covered this week."
    }
    merged_nodes[week_root_id] = week_root

    for g in graphs_list:
        local_nodes = {n['id']: n for n in g.get('nodes', [])}
        local_edges = g.get('edges', [])

        # 2. Find File Root
        file_root_id = None
        for nid, node in local_nodes.items():
            if node.get('isRoot') or node.get('level') == 0:
                file_root_id = nid
                break

        # 3. BFS to shift levels (file root → level 1)
        if file_root_id:
            queue = deque([(file_root_id, 1)])
            visited = {file_root_id}
            adj = defaultdict(list)
            for e in local_edges:
                adj[e['source']].append(e['target'])
            while queue:
                curr, new_lvl = queue.popleft()
                if curr in local_nodes:
                    local_nodes[curr]['level'] = new_lvl
                    local_nodes[curr]['isRoot'] = False
                for child in adj[curr]:
                    if child not in visited:
                        visited.add(child)
                        queue.append((child, new_lvl + 1))

        # 4. Merge Nodes with label-similarity dedup
        for nid, node in local_nodes.items():
            label = node.get('label', nid)

            # Check if a node with a similar label already exists
            canonical_id = None
            for existing_id, existing_node in merged_nodes.items():
                if existing_id == week_root_id:
                    continue
                if _similar_labels(label, existing_node.get('label', existing_id)):
                    canonical_id = existing_id
                    break

            if canonical_id:
                # Duplicate found — remap this ID to the canonical one
                id_remap[nid] = canonical_id
                # Merge slide_nums from the duplicate into the canonical node
                existing = merged_nodes[canonical_id]
                existing_nums = existing.get('slide_nums', [])
                new_nums = node.get('slide_nums', [])
                if isinstance(new_nums, list) and isinstance(existing_nums, list):
                    existing['slide_nums'] = sorted(set(existing_nums) | set(new_nums))
            else:
                # New unique concept
                id_remap[nid] = nid
                merged_nodes[nid] = node

        # 5. Merge Edges (rewrite IDs through remap table)
        for edge in local_edges:
            src = id_remap.get(edge['source'], edge['source'])
            tgt = id_remap.get(edge['target'], edge['target'])
            sig = (src, tgt)
            if sig not in seen_edges and src != tgt:
                seen_edges.add(sig)
                merged_edges.append({"source": src, "target": tgt, "relation": edge.get('relation', 'related')})

        # 6. Connect File Root → Week Root
        canonical_file_root = id_remap.get(file_root_id, file_root_id) if file_root_id else None
        if canonical_file_root and (week_root_id, canonical_file_root) not in seen_edges:
            merged_edges.append({"source": week_root_id, "target": canonical_file_root, "relation": "topic"})
            seen_edges.add((week_root_id, canonical_file_root))

    return {"nodes": list(merged_nodes.values()), "edges": merged_edges}


def extract_mcqs_from_response(resp):
    """
    Handles LLM output: accepts string or list, safely strips code block fencing, parses JSON,
    and recovers MCQs robustly.
    """
    # 1. If resp is already parsed (list or dict), return MCQs
    if isinstance(resp, list):
        return resp
    if isinstance(resp, dict) and "mcqs" in resp:
        return resp["mcqs"]

    # 2. If resp is a string, clean it up
    if isinstance(resp, str):
        text = resp.strip()
        # Remove markdown code fence if present
        if text.startswith("```"):
            lines = text.splitlines()
            # Remove opening ```, possibly with 'json'
            if lines and lines[0].strip().startswith("```"):
                lines = lines[1:]
            # Remove closing ```
            if lines and lines[-1].strip().startswith("```"):
                lines = lines[:-1]
            text = "\n".join(lines).strip()
        try:
            items = json.loads(text)
            if isinstance(items, dict) and "mcqs" in items:
                return items["mcqs"]
            if isinstance(items, list):
                return items
            return []
        except Exception as e:
            print("[MCQ PARSE ERROR]", e)
            return []
    print("[MCQ PARSE ERROR] Response is neither str nor list nor dict:", type(resp))
    return []

def pick_next_concept(concept_ids, responses, graph_nodes_map, include_spaced=True):
    """
    Intelligent Selection based on Type & Hierarchy.
    
    Priority Order:
    1. Unseen 'structure' or 'root' concepts (Foundational)
    2. Unseen 'algorithm' or 'concept' (Core material)
    3. Unseen 'detail' or 'example' (Nuance)
    4. Review (Spaced Repetition)
    """
    # 1. Identify what's seen vs unseen
    seen_ids = {r.get("concept_id") for r in responses.values() if r.get("concept_id")}
    unseen_ids = [c for c in concept_ids if c not in seen_ids]

    # 2. Define Type Weights (Lower number = Higher Priority)
    type_priority = {
        "root": 0,
        "structure": 1,
        "algorithm": 1,
        "concept": 2,
        "application": 3,
        "detail": 4,
        "example": 5,
        "trivia": 6
    }

    # 3. Helper to get priority score
    def get_score(cid):
        node = graph_nodes_map.get(cid, {})
        # Use 'type' if available, otherwise fallback to 'concept'
        ctype = node.get("type", "concept").lower()
        # Add a tiny random float (0.0-0.9) to shuffle items within the same tier
        return type_priority.get(ctype, 3) + random.random()

    # 4. Pick Best Unseen
    if unseen_ids:
        # Sort by priority score (ascending)
        unseen_ids.sort(key=get_score)
        return unseen_ids[0]

    # 5. Review Mode (If all unseen are done)
    if include_spaced and seen_ids:
        # 30% chance to pick a harder concept to review, 70% random
        if random.random() < 0.3:
            # Pick a 'structure' or 'algorithm' to review
            high_value_seen = [cid for cid in seen_ids if graph_nodes_map.get(cid, {}).get('type') in ['structure', 'algorithm']]
            if high_value_seen:
                return random.choice(high_value_seen)
        
        return random.choice(list(seen_ids))

    return random.choice(concept_ids)


# ==== LECTURE UPLOAD, PROCESSING, AND STATUS ====

@app.post("/api/upload-lecture/")
async def upload_lecture(
    background_tasks: BackgroundTasks, # Import BackgroundTasks from fastapi
    file: UploadFile = File(...),
    course_id: str = Form(...),
    week: int = Form(...),
    db: Session = Depends(get_db)
):
    # 1. Save File
    file_id = str(uuid.uuid4())
    safe_filename = f"{file_id}_{file.filename}"
    file_path = os.path.join(UPLOAD_DIR, safe_filename)
    
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    # 2. Create DB Record
    processing_record = LectureProcessing(
        id=file_id,
        course_id=course_id,
        week=week,
        file_name=file.filename,
        status="pending",
        progress=0
    )
    db.add(processing_record)
    db.commit()

    # 3. Trigger Background Processing (Calls your GOOD function)
    background_tasks.add_task(
        process_lecture_and_kg, # The big function above
        file_path,
        file_id, 
        course_id, 
        week, 
        file.filename, 
        file_id  # using file_id as processing_id for simplicity
    )

    return {"message": "Upload started", "processing_id": file_id}


@app.get("/api/lecture-status/")
async def lecture_status(processing_id: str, db: Session = Depends(get_db)):
    job = db.query(LectureProcessing).filter(LectureProcessing.id == processing_id).first()
    if not job:
        raise HTTPException(status_code=404, detail="Job not found")
    return {
        "status": job.status,
        "progress": job.progress,
        "error": job.error_message,
        "file_name": job.file_name
    }

@app.get("/api/lecture-history/")
async def lecture_history(course_id: str, db: Session = Depends(get_db)):
    jobs = db.query(LectureProcessing).filter(LectureProcessing.course_id == course_id).all()
    return [{
        "id": j.id, "week": j.week, "file_name": j.file_name,
        "status": j.status, "progress": j.progress, "error": j.error_message
    } for j in jobs]

@app.get("/api/segments/")
async def get_segments(course_id: str, week: int, db: Session = Depends(get_db)):
    segs = db.query(Segment).filter_by(course_id=course_id, week=week).order_by(Segment.segment_index).all()
    return [
        {"id": s.id, "segment_index": s.segment_index, "title": s.title, "keywords": s.keywords}
        for s in segs
    ]

@app.get("/api/segment/{segment_id}")
async def get_segment_detail(segment_id: str, db: Session = Depends(get_db)):
    s = db.query(Segment).filter_by(id=segment_id).first()
    if not s:
        raise HTTPException(status_code=404, detail="Segment not found")
    return {
        "id": s.id,
        "course_id": s.course_id,            # <-- add this
        "week": s.week,                      # <-- add this
        "segment_index": s.segment_index,    # <-- and this
        "title": s.title,
        "content": s.content,
        "summary": s.summary,
        "keywords": s.keywords
    }

@app.get("/api/knowledge-graph/list")
def list_knowledge_graphs(
    course_id: str, 
    week: int, 
    db: Session = Depends(get_db)
):
    """
    Returns a list of all available knowledge graphs for this week.
    Example: [{ type: 'master', file: null }, { type: 'file', file: 'Week1.pdf' }]
    """
    graphs = db.query(KnowledgeGraph.graph_type, KnowledgeGraph.source_file).filter(
        KnowledgeGraph.course_id == course_id,
        KnowledgeGraph.week == week
    ).all()
    
    results = []
    # Always ensure 'master' is an option implicitly, but good to check if it exists
    has_master = False
    
    for g in graphs:
        if g.graph_type == 'master':
            has_master = True
        elif g.graph_type == 'file':
            results.append({
                "id": g.source_file, # Use filename as ID for simplicity
                "name": g.source_file,
                "type": "file"
            })
            
    # Always prepend Master
    # (or only if has_master is True, but usually you always want the option)
    results.insert(0, { "id": "master", "name": "Combined Master Graph", "type": "master" })
    
    return results


@app.get("/api/knowledge-graph")
def get_knowledge_graph(
    courseid: str, 
    week: int, 
    graph_type: str = "master", # Default to master
    source_file: str = None,
    db: Session = Depends(get_db)
):
    query = db.query(KnowledgeGraph).filter(
        KnowledgeGraph.course_id == courseid,
        KnowledgeGraph.week == week
    )
    
    if source_file:
        query = query.filter(KnowledgeGraph.source_file == source_file, KnowledgeGraph.graph_type == "file")
    else:
        query = query.filter(KnowledgeGraph.graph_type == graph_type)
        
    kg = query.first()
    
    if not kg:
        return {"nodes": [], "edges": []}
        
    return {
        "nodes": json.loads(kg.node_data),
        "edges": json.loads(kg.edge_data)
    }


# ==== RETRIEVAL PRACTICE: MCQ GENERATION (basic stub) ====
@app.post("/api/generate-mcqs/")
async def generate_mcqs(payload: MCQConceptModel = Body(...)):
    # Compose a targeted LLM prompt
    prompt = (
        f"You are a teaching assistant. Create several relevant multiple-choice questions about the concept '{payload.concept_id}' "
        "based only on its context in the following lecture notes. Use the summary for focus and the detailed contents for depth. "
        "Make sure all questions focus specifically on this concept as covered in the summary/content below. "
        "The number of questions should depend on its importance and how much material is present. "
        "Make sure EACH QUESTION is accompanied by exactly 4 numbered options, and provide the correct answer as an 'answer' field (the text matching one of the options, not the letter/number). "
        "Respond ONLY as a JSON list in the format:"
        '[{"question": "...", "options": ["A", "B", "C", "D"], "answer": "..."}]\n'
        f"\nConcept summary:\n{payload.summary[:600]}"
        f"\n\nRelated contents (for MCQ details):\n{payload.contents[:1200]}"
    )
    try:
        model = genai.GenerativeModel("models/gemini-2.5-flash")
        response = model.generate_content(prompt)
        model_output = str(getattr(response, "text", getattr(response, "candidates", response)))
        mcqs = extract_mcqs_from_response(model_output)
        out = []
        for item in mcqs:
            if isinstance(item, dict) and all(k in item for k in ("question", "options", "answer")) and len(item["options"]) == 4:
                out.append(item)
        return {"mcqs": out if out else [{"question": f"No MCQs found for concept '{payload.concept_id}'.", "options": [], "answer": ""}]}
    except Exception as ex:
        print("[MCQ GENERATION ERROR]", ex)
        return {"mcqs": [{"question": "Sample fallback MCQ: LLM Error, please try again.", "options": ["A", "B", "C", "D"], "answer": "A"}]}
    except Exception as ex:
        print("[MCQ GENERATION ERROR]", ex)
        return {"mcqs": [{
            "question": "Sample fallback MCQ: LLM Error, please try again.",
            "options": ["A", "B", "C", "D"],
            "answer": "A"
        }]}


@app.post("/api/generate-mcqs-kg")
async def generate_mcqs_kg(payload: dict = Body(...), db: Session = Depends(get_db)):
    course_id = payload.get("course_id")
    week = payload.get("week")
    concept_id = payload.get("concept_id") or payload.get("segment_id")
    quiz_id = payload.get("quiz_id")
    
    print(f"[MCQ GEN] Starting generation for concept: {concept_id}, quiz: {quiz_id}")
    
    # 1. Determine Difficulty Range (NOT affected by settings - always generate all)
    all_difficulties = ["Easy", "Medium", "Hard"]
    
    # 2. Determine Bloom Level Range (NOT affected by settings - always generate all)
    bloom_hierarchy = ["Remember", "Understand", "Apply", "Analyze", "Evaluate", "Create"]
    
    # ✅ GENERATION should ALWAYS create all Bloom levels and difficulties
    # ✅ FILTERING happens later in student quiz start endpoint
    generation_bloom_levels = bloom_hierarchy[:]
    generation_bloom_str = ", ".join([f"'{x}'" for x in generation_bloom_levels])
    
    print(f"[MCQ GEN] Will generate across ALL Bloom levels: {generation_bloom_str}")
    print(f"[MCQ GEN] Will generate across ALL difficulties: {all_difficulties}")

    # 3. Fetch KG Context (Robust Attribute Access)
    kg_entry = None
    try:
        kg_entry = db.query(KnowledgeGraph).filter(
            KnowledgeGraph.course_id == course_id,
            KnowledgeGraph.week == week,
            KnowledgeGraph.graph_type == "master"  # Use the master merged graph
        ).first()
    except Exception as e:
        print(f"[MCQ GEN] KG Lookup Error: {e}")

    selected_summary = ""
    selected_contents = ""

    if kg_entry:
        try:
            raw_node_data = kg_entry.node_data  # Use snake_case from unified models
            
            if raw_node_data:
                kg_nodes = json.loads(raw_node_data)
                selected_node = next((n for n in kg_nodes if n.get('id') == concept_id), None)
                if selected_node:
                    selected_summary = selected_node.get('summary', '')
                    selected_contents = selected_node.get('contents', selected_summary)
            else:
                print("[MCQ GEN] Warning: KG entry found but no node data.")
                
        except json.JSONDecodeError:
            print("[MCQ GEN] Error decoding KG JSON")
        except Exception as e:
            print(f"[MCQ GEN] Error parsing KG nodes: {e}")

    # 4. ✅ Build Bloom Instruction (GENERATE ALL LEVELS - no constraints)
    bloom_instruction = f"""
    BLOOM'S TAXONOMY REQUIREMENT:
    - Generate questions across ALL cognitive levels: {generation_bloom_str}
    - Bloom hierarchy: Remember < Understand < Apply < Analyze < Evaluate < Create
    - Create a DIVERSE mix of questions at different cognitive levels
    - Include at least one question for each Bloom level if possible
    - IMPORTANT: You MUST include the 'bloom_level' field for every question
    """

    # 5. Prompt Generation
    prompt = f"""
As an expert computer science instructor, create high-quality multiple-choice questions (MCQs) for the concept "{concept_id}".

TARGET DIFFICULTIES: Generate questions at Easy, Medium, and Hard levels (mix them)

{bloom_instruction}

STRICT FORMATTING RULES:
1. Questions must be stand-alone and not reference "the text" or "according to..."
2. Provide exactly 4 options per question (as a list)
3. REQUIRED: Include 'difficulty' field with one of: "Easy", "Medium", "Hard"
4. REQUIRED: Include 'bloom_level' field with one of: {generation_bloom_str}
5. Respond ONLY as a valid JSON array (no markdown, no backticks, no code blocks)

Response Format (EXACT JSON):
[
  {{
    "question": "What is the time complexity of binary search?",
    "options": ["O(1)", "O(log n)", "O(n)", "O(n²)"],
    "answer": "O(log n)",
    "difficulty": "Medium",
    "bloom_level": "Remember"
  }},
  {{
    "question": "How would you modify binary search to find the first occurrence of a duplicate element?",
    "options": ["Continue searching left after finding target", "Return immediately", "Use linear search", "Sort the array first"],
    "answer": "Continue searching left after finding target",
    "difficulty": "Hard",
    "bloom_level": "Apply"
  }}
]

--- CONCEPT CONTEXT ---
Summary: {selected_summary[:800] if selected_summary else "No summary available."}
Details: {selected_contents[:1200] if selected_contents else "No additional details."}

Generate 5-7 diverse questions covering different Bloom levels and difficulties.
"""

    try:
        model = genai.GenerativeModel('models/gemini-2.5-flash')
        response = model.generate_content(prompt)
        model_output = getattr(response, 'text', str(response))
        
        print(f"[MCQ GEN] Raw LLM output length: {len(model_output)}")
        
        mcqs = extract_mcqs_from_response(model_output)
        
        print(f"[MCQ GEN] Extracted {len(mcqs)} raw MCQs from response")

        out = []
        
        for idx, item in enumerate(mcqs):
            if isinstance(item, dict) and 'question' in item and 'options' in item and 'answer' in item:
                
                # ✅ Ensure difficulty exists and is valid
                if 'difficulty' not in item or item['difficulty'] not in all_difficulties:
                    item['difficulty'] = random.choice(all_difficulties)
                    print(f"[MCQ GEN] MCQ {idx}: Set default difficulty -> {item['difficulty']}")
                
                # ✅ Ensure bloom_level exists and is valid
                if 'bloom_level' not in item or item['bloom_level'] not in bloom_hierarchy:
                    item['bloom_level'] = random.choice(bloom_hierarchy)
                    print(f"[MCQ GEN] MCQ {idx}: Set default bloom_level -> {item['bloom_level']}")
                else:
                    print(f"[MCQ GEN] MCQ {idx}: bloom_level = {item['bloom_level']}, difficulty = {item['difficulty']}")
                
                # Check for hallucinations (questions that reference "the text")
                q_text = item['question'].lower()
                if any(phrase in q_text for phrase in ["provided text", "according to", "the text states", "in the passage", "based on the passage"]):
                    print(f"[MCQ GEN] Skipping hallucinated question: {item['question'][:60]}...")
                    continue

                # ✅ Attach concept_id
                item['concept_id'] = concept_id
                out.append(item)
            else:
                print(f"[MCQ GEN] Skipping malformed MCQ at index {idx}: {item}")
        
        # ✅ Log distribution
        bloom_distribution = {}
        difficulty_distribution = {}
        for mcq in out:
            bloom_distribution[mcq['bloom_level']] = bloom_distribution.get(mcq['bloom_level'], 0) + 1
            difficulty_distribution[mcq['difficulty']] = difficulty_distribution.get(mcq['difficulty'], 0) + 1
        
        print(f"[MCQ GEN] ✅ Generated {len(out)} valid MCQs")
        print(f"[MCQ GEN] Bloom distribution: {bloom_distribution}")
        print(f"[MCQ GEN] Difficulty distribution: {difficulty_distribution}")
        
        return {"mcqs": out}

    except Exception as ex:
        print(f"[MCQ GEN] ❌ GENERATION ERROR: {ex}")
        traceback.print_exc()
        return {"mcqs": []}




@app.get("/api/mcqs/")
async def get_mcqs(segment_id: str, db: Session = Depends(get_db)):
    # Return MCQs for the segment
    mcqs = db.query(MCQ).filter_by(segment_id=segment_id).all()
    return [
        {"question": m.question, "options": m.options, "answer": m.answer}
        for m in mcqs
    ]

@app.post("/api/auth/signup")
async def signup(email: str = Form(...), password: str = Form(...), role: str = Form(...), db: Session = Depends(get_db)):
    # Check if already exists
    user = db.query(User).filter(User.email == email).first()
    if user:
        raise HTTPException(status_code=409, detail="Email already registered")
    hashed_pw = hashlib.sha256(password.encode()).hexdigest()
    new_user = User(id=str(uuid.uuid4()), email=email, password=hashed_pw, role=role)
    db.add(new_user)
    db.commit()
    db.refresh(new_user)
    return {"id": new_user.id, "email": new_user.email, "role": new_user.role}

@app.post("/api/auth/login")
async def login(email: str = Form(...), password: str = Form(...), db: Session = Depends(get_db)):
    hashed_pw = hashlib.sha256(password.encode()).hexdigest()
    user = db.query(User).filter(User.email == email, User.password == hashed_pw).first()
    if not user:
        raise HTTPException(status_code=401, detail="Invalid email or password")
    return {"id": user.id, "email": user.email, "role": user.role}

# ==== COURSE MANAGEMENT ====
@app.get("/api/courses")
async def get_courses(instructor_id: str = None, db: Session = Depends(get_db)):
    if instructor_id:
        courses = db.query(Course).filter(Course.instructor_id == instructor_id).all()
    else:
        courses = db.query(Course).all()
    return [
        {
            "id": c.id,
            "name": c.name,
            "instructor_id": c.instructor_id,
            "created_at": c.created_at.isoformat()
        }
        for c in courses
    ]

@app.post("/api/courses")
async def create_course(name: str = Form(...), instructor_id: str = Form(...), db: Session = Depends(get_db)):
    new_course = Course(id=str(uuid.uuid4()), name=name, instructor_id=instructor_id)
    db.add(new_course)
    db.commit()
    db.refresh(new_course)
    return {
        "id": new_course.id,
        "name": new_course.name,
        "instructor_id": new_course.instructor_id,
        "created_at": new_course.created_at.isoformat()
    }

# ==== ENROLLMENT: Students enroll in courses ====
@app.post("/api/enroll")
async def enroll(course_id: str = Form(...), student_id: str = Form(...), db: Session = Depends(get_db)):
    # assumes Enrollment model exists with id, course_id, student_id
    enrollment = db.query(Enrollment).filter(
        Enrollment.course_id == course_id,
        Enrollment.student_id == student_id
    ).first()
    if enrollment:
        raise HTTPException(status_code=409, detail="Already enrolled")
    new_enrollment = Enrollment(id=str(uuid.uuid4()), course_id=course_id, student_id=student_id)
    db.add(new_enrollment)
    db.commit()
    db.refresh(new_enrollment)
    return {"status": "enrolled", "enrollment_id": new_enrollment.id}

@app.get("/api/student-courses")
async def get_student_courses(student_id: str, db: Session = Depends(get_db)):
    enrollments = db.query(Enrollment).filter(Enrollment.student_id == student_id).all()
    course_ids = [en.course_id for en in enrollments]
    courses = db.query(Course).filter(Course.id.in_(course_ids)).all()
    return [
        {
            "id": c.id,
            "name": c.name,
            "instructor_id": c.instructor_id,
            "created_at": c.created_at.isoformat()
        }
        for c in courses
    ]

# ==== MODULES & WEEKS (for week-based structure) ====
@app.get("/api/modules")
async def get_modules(course_id: str, db: Session = Depends(get_db)):
    # assumes Module model with id, name, course_id, week
    modules = db.query(Module).filter(Module.course_id == course_id).order_by(Module.week).all()
    return [
        {
            "id": m.id,
            "name": m.name,
            "week": m.week,
            "course_id": m.course_id
        }
        for m in modules
    ]

@app.post("/api/modules")
async def create_module(name: str = Form(...), course_id: str = Form(...), week: int = Form(...), db: Session = Depends(get_db)):
    # New teaching module for a given week in course
    mod = Module(id=str(uuid.uuid4()), name=name, course_id=course_id, week=week)
    db.add(mod)
    db.commit()
    db.refresh(mod)
    return {
        "id": mod.id,
        "name": mod.name,
        "week": mod.week,
        "course_id": mod.course_id
    }

@app.get("/api/debug/segment/{segment_id}")
async def debug_segment(segment_id: str, db: Session = Depends(get_db)):
    s = db.query(Segment).filter_by(id=segment_id).first()
    if not s:
        raise HTTPException(status_code=404, detail="Segment not found")
    return {
        "id": s.id,
        "title": s.title,
        "content": s.content[:200] + "..." if s.content else None,
        "summary": s.summary,
        "keywords": s.keywords,
        "summary_length": len(s.summary) if s.summary else 0,
        "has_summary": bool(s.summary and s.summary.strip())
    }

@app.delete("/api/upload/{upload_id}")
def delete_upload(upload_id: str, db: Session = Depends(get_db)):
    record = db.query(LectureProcessing).filter(LectureProcessing.id == upload_id).first()
    if not record:
        raise HTTPException(status_code=404, detail="Upload not found")
    
    # DELETE GRAPH LOGIC
    db.query(KnowledgeGraph).filter(
        KnowledgeGraph.course_id == record.course_id,
        KnowledgeGraph.week == record.week
    ).delete()
    
    db.delete(record)
    db.commit()
    return {"status": "deleted"}

@app.delete("/api/quiz/{quiz_id}")
async def delete_quiz(quiz_id: str, db: Session = Depends(get_db)):
    """
    Deletes a quiz and all its related data (MCQs, attempts, settings).
    """
    # 1. Find the quiz
    quiz = db.query(Quiz).filter(Quiz.id == quiz_id).first()
    if not quiz:
        raise HTTPException(status_code=404, detail="Quiz not found")

    # 2. Delete related MCQs
    db.query(MCQ).filter(MCQ.quiz_id == quiz_id).delete()
    
    # 3. Delete related Attempts (optional, but good for cleanup)
    db.query(QuizAttempt).filter(QuizAttempt.quiz_id == quiz_id).delete()
    
    # 4. Delete Settings
    # Try both attribute names for safety based on our previous fixes
    try:
        db.query(QuizSettings).filter(QuizSettings.quiz_id == quiz_id).delete()
    except:
        db.query(QuizSettings).filter(QuizSettings.quizid == quiz_id).delete()

    # 5. Delete the Quiz itself
    db.delete(quiz)
    
    try:
        db.commit()
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Failed to delete quiz: {e}")

    return {"status": "deleted", "quiz_id": quiz_id}

# Add WebSocket endpoint
@app.websocket("/ws/{course_id}/{week}")
async def websocket_endpoint(websocket: WebSocket, course_id: str, week: int):
    await manager.connect(websocket)
    try:
        while True:
            data = await websocket.receive_text()
            await manager.send_personal_message(f"Message: {data}", websocket)
    except WebSocketDisconnect:
        manager.disconnect(websocket)

# ✅ WebSocket endpoint for instructors 
@app.websocket("/ws/quiz/{quiz_id}")
async def websocket_quiz_endpoint(websocket: WebSocket, quiz_id: str):
    await websocket.accept()
    
    # ✅ Allow connections without auth for now
    user_id = "instructor"  # or websocket.query_params.get("user_id")
    
    await manager.connect(quiz_id, user_id, websocket)
    try:
        while True:
            data = await websocket.receive_text()
            # handle messages
    except WebSocketDisconnect:
        manager.disconnect(quiz_id, user_id)
        await websocket.close()


# --- 1. Create a quiz (Instructor) ---
@app.post("/api/quiz/create")
async def create_quiz(request: QuizCreate, db: Session = Depends(get_db)):
    print("✅ Creating quiz:", request.course_id, len(request.concept_ids))
    
    # 1. CREATE QUIZ ONLY
    qid = str(uuid.uuid4())
    quiz = Quiz(
        id=qid,
        name=request.name,
        course_id=request.course_id,
        week=request.week,
        concept_ids=request.concept_ids,
        instructor_id=request.instructor_id
    )
    db.add(quiz)
    db.commit()
    
    print(f"📋 Quiz {qid} created with {len(request.concept_ids)} concepts")
    return {"quiz_id": qid, "message": "Quiz created, ready for MCQ generation"}


# --- 2. List all quizzes for a course/week (for student or dashboard) ---
@app.get("/api/quiz/list")
async def list_quizzes(course_id: str, week: int, db: Session = Depends(get_db)):
    quizzes = db.query(Quiz).filter_by(course_id=course_id, week=week).all()
    return [
        {"id": q.id, "name": q.name, "concept_ids": q.concept_ids}
        for q in quizzes
    ]

# --- 3. Start or resume a quiz attempt (one per student per quiz) ---
@app.post("/api/quiz/attempt/start")
async def start_quiz_attempt(
    quiz_id: str = Form(...),
    student_id: str = Form(...),
    db: Session = Depends(get_db)
):
    """Start or resume a quiz attempt for a student."""
    attempt = db.query(QuizAttempt).filter_by(quiz_id=quiz_id, student_id=student_id).first()
    if attempt:
        return {"attempt_id": attempt.id, "resume": True}

    new_attempt = QuizAttempt(
        id=str(uuid.uuid4()),
        quiz_id=quiz_id,
        student_id=student_id,
        responses={}
    )
    db.add(new_attempt)
    db.commit()
    db.refresh(new_attempt)
    return {"attempt_id": new_attempt.id, "resume": False}

# --- 4. Fetch next adaptive MCQ (students; Gemini/KG) ---
@app.get("/api/quiz/attempt/next")
async def get_next_mcq_gemini(attempt_id: str, db: Session = Depends(get_db)):
    attempt = db.query(QuizAttempt).filter_by(id=attempt_id).first()
    if not attempt:
        raise HTTPException(status_code=404, detail="Attempt not found")

    quiz = db.query(Quiz).filter_by(id=attempt.quiz_id).first()
    settings = db.query(QuizSettings).filter_by(quiz_id=quiz.id).first()

    # 1. Check Limits
    max_questions = settings.max_questions if settings else 10
    responses = attempt.responses or {}
    if len(responses) >= max_questions:
        return {"done": True, "reason": "max_questions_reached"}

    # 2. Load Graph Nodes for Intelligent Selection
    kg_nodes_map = {} # Default to empty dict
    
    try:
        kg_row = db.execute(
            sql_text("SELECT node_data FROM knowledge_graph WHERE course_id=:c AND week=:w"),
            {"c": quiz.course_id, "w": quiz.week}
        ).fetchone()

        if kg_row and kg_row[0]:
            nodes = json.loads(kg_row[0])
            kg_nodes_map = {n['id']: n for n in nodes}
    except Exception as e:
        print(f"[WARNING] Graph load failed for Quiz {quiz.id}: {e}")
        # Proceed with empty map - logic will fallback to random choice
    # 3. Pick Concept
    concept_ids = quiz.concept_ids or []
    next_concept_id = pick_next_concept(
        concept_ids, 
        responses, 
        kg_nodes_map, # <-- Now passing the map!
        include_spaced=(settings.include_spaced if settings else False)
    )

    # 4. Try Cached MCQs First
    existing_mcqs = db.query(MCQ).filter(
        MCQ.quiz_id == quiz.id, 
        MCQ.concept_id == next_concept_id
    ).all()
    
    # Filter out already answered
    answered_ids = set(responses.keys())
    valid_mcqs = [m for m in existing_mcqs if m.id not in answered_ids]

    if valid_mcqs:
        # Pick one (random or by difficulty if you implemented that)
        selected = random.choice(valid_mcqs)
        return {
            "mcq_id": selected.id,
            "question": selected.question,
            "options": selected.options,
            "difficulty": selected.difficulty,
            "bloom_level": selected.bloom_level, 
            "concept_id": selected.concept_id
        }

    # 5. Generate New if None Exist
    print(f"⚡ Generating fresh MCQ for {next_concept_id}...")
    kg_payload = {
        "mcq_id": new_mcq.id,
        "question": new_mcq.question,
        "options": new_mcq.options,
        "concept_id": new_mcq.concept_id,
        "difficulty": new_mcq.difficulty, 
        "bloom_level": new_mcq.bloom_level
    }
    
    try:
        gen_resp = await generate_mcqs_kg(kg_payload, db)
        new_mcqs_data = gen_resp.get("mcqs", [])
        
        if new_mcqs_data:
            # SAVE IT IMMEDIATELY so it has a real ID
            first_q = new_mcqs_data[0]
            new_mcq = MCQ(
                id=str(uuid.uuid4()),
                quiz_id=quiz.id,
                concept_id=next_concept_id,
                question=first_q["question"],
                options=first_q["options"],
                answer=first_q["answer"],
                difficulty=first_q.get("difficulty", "Medium"),
                bloom_level=item["bloom_level"]
            )
            db.add(new_mcq)
            db.commit()
            
            return {
                "mcq_id": new_mcq.id,
                "question": new_mcq.question,
                "options": new_mcq.options,
                "concept_id": new_mcq.concept_id
            }
            
    except Exception as e:
        print(f"Gen Failed: {e}")
    
    return {"done": True, "reason": "generation_failed"}

@app.post("/api/quiz/{quiz_id}/mcqs")  # or whatever
async def add_generated_mcqs(quiz_id: str, payload: dict, db):
    mcqs_data = payload["mcqs"]
    for item in mcqs_data:
        mcq = MCQ(
            id=str(uuid.uuid4()),
            quiz_id=quiz_id,
            question=item["question"],
            options=json.dumps(item["options"]),
            answer=item["answer"],
            concept_id=item["concept_id"],
            difficulty=item["difficulty"],
            bloom_level=item["bloom_level"],  # ✅ ADD THIS LINE
        )
        db.add(mcq)
    db.commit()
    return {"added": len(mcqs_data)}


# --- 5. Submit MCQ answer for attempt (DB or Gemini MCQ) ---
@app.post("/api/quiz/attempt/submit")
async def submit_mcq_answer(
    attempt_id: str = Form(...),
    mcq_id: str = Form(...),
    selected: str = Form(...),
    answer: str = Form(None),          # for Gemini MCQs
    concept_id: str = Form(None),      # NEW
    db: Session = Depends(get_db)
):
    attempt = db.query(QuizAttempt).filter_by(id=attempt_id).first()
    if not attempt:
        raise HTTPException(status_code=404, detail="Attempt not found")

    mcq = db.query(MCQ).filter_by(id=mcq_id).first()

    if mcq is not None:
        correct = selected.strip() == mcq.answer.strip()
        answer_val = mcq.answer
        question = mcq.question
        concept_id_used = mcq.concept_id or concept_id
    elif answer is not None:
        correct = selected.strip() == answer.strip()
        answer_val = answer
        question = "(Generated MCQ)"
        concept_id_used = concept_id
    else:
        raise HTTPException(status_code=400, detail="No answer key for this MCQ")

    responses = attempt.responses or {}
    responses[mcq_id] = {
        "selected": selected,
        "correct": correct,
        "concept_id": concept_id_used,
    }
    attempt.responses = responses
    db.commit()

    mcq_resp = MCQResponse(
        id=str(uuid.uuid4()),
        attempt_id=attempt_id,
        mcq_id=mcq_id,
        question=question,
        selected=selected,
        correct=correct,
    )
    db.add(mcq_resp)
    db.commit()

    return {"correct": correct, "answer": answer_val}

# --- 6. Get summary of quiz attempt ---
@app.get("/api/quiz/attempt/state")
async def quiz_attempt_state(attempt_id: str, db: Session = Depends(get_db)):
    """Return a summary (attempted/correct count, responses) for an attempt."""
    attempt = db.query(QuizAttempt).filter_by(id=attempt_id).first()
    if not attempt:
        raise HTTPException(status_code=404, detail="Attempt not found")
    responses = attempt.responses or {}
    return {
        "attempted": len(responses),
        "correct": sum(1 for r in responses.values() if r.get("correct")),
        "responses": responses
    }

# --- 7. Quiz aggregate stats (by instructor) ---
@app.get("/api/quiz/stats")
async def quiz_stats(quiz_id: str, db: Session = Depends(get_db)):
    # Aggregate stats for a quiz
    attempts = db.query(QuizAttempt).filter_by(quiz_id=quiz_id).all()
    summary = []
    
    for att in attempts:
        # 1. Handle if 'responses' is None
        raw_responses = att.responses or {}
        
        # 2. Handle if 'responses' is a JSON string (SQLite/SQLAlchemy quirk)
        if isinstance(raw_responses, str):
            try:
                responses = json.loads(raw_responses)
            except:
                responses = {}
        else:
            responses = raw_responses

        # 3. Calculate Correct Count Safely
        correct_count = 0
        if isinstance(responses, dict):
             for v in responses.values():
                 # Handle if the value inside the dict is ALSO a string (double encoded)
                 if isinstance(v, str):
                     try:
                         v = json.loads(v)
                     except:
                         continue # Skip malformed data
                 
                 if isinstance(v, dict) and v.get("correct"):
                     correct_count += 1
        
        summary.append({
            "attempt_id": att.id,
            "student_id": att.student_id,
            "attempted": len(responses) if isinstance(responses, dict) else 0,
            "correct": correct_count,
            "score": att.score
        })

    return {
        "total_attempts": len(attempts),
        "stats": summary
    }

# --- 8. Home/root route ---
@app.get("/")
async def root():
    return {"message": "AILA Backend (SQLite) is running"}

# --- 9. Delete upload/file endpoint (optional housekeeping) ---
@app.post("/api/delete-upload")
async def delete_upload(uploadid: str = Form(...), db: Session = Depends(get_db)):
    """Delete a previously uploaded lecture file."""
    upload = db.query(LectureUpload).filter(LectureUpload.id == uploadid).first()
    if not upload:
        raise HTTPException(status_code=404, detail="File not found")
    if upload.fileurl and os.path.exists(upload.fileurl):
        os.remove(upload.fileurl)
    db.delete(upload)
    db.commit()
    return {"status": "deleted"}

@app.get("/api/quiz/settings/{quiz_id}", response_model=QuizSettingsOut)
async def get_quiz_settings(quiz_id: str, db: Session = Depends(get_db)):
    qs = None
    try:
        qs = db.query(QuizSettings).filter(QuizSettings.quiz_id == quiz_id).first()
    except AttributeError:
        qs = db.query(QuizSettings).filter(QuizSettings.quizid == quiz_id).first()
    
    if not qs:
        quiz = db.query(Quiz).filter(Quiz.id == quiz_id).first()
        default_week = quiz.week if quiz else 1
        return {
            "id": -1,
            "quizid": quiz_id,
            "quiz_id": quiz_id,
            "week": default_week,
            "mindifficulty": "Easy",
            "maxdifficulty": "Hard",
            "min_bloom_level": "Remember", 
            "max_bloom_level": "Create",    
            "maxquestions": 10,
            "allowedretries": 3,
            "feedbackstyle": "Immediate",
            "includespaced": False
        }

    # Get existing values with safe attribute access
    return {
        "id": qs.id,
        "quizid": getattr(qs, 'quiz_id', quiz_id),
        "quiz_id": getattr(qs, 'quiz_id', quiz_id),
        "week": qs.week,
        "mindifficulty": getattr(qs, 'min_difficulty', 'Easy'),
        "maxdifficulty": getattr(qs, 'max_difficulty', 'Hard'),
        "min_bloom_level": getattr(qs, 'min_bloom_level', 'Remember'),  
        "max_bloom_level": getattr(qs, 'max_bloom_level', 'Create'),    
        "maxquestions": getattr(qs, 'max_questions', 10),
        "allowedretries": getattr(qs, 'allowed_retries', 3),
        "feedbackstyle": getattr(qs, 'feedback_style', 'Immediate'),
        "includespaced": getattr(qs, 'include_spaced', False)
    }


@app.post("/api/quiz/settings/{quiz_id}", response_model=QuizSettingsOut)
async def upsert_quiz_settings(quiz_id: str, payload: QuizSettingsIn, db: Session = Depends(get_db)):
    qs = db.query(QuizSettings).filter(QuizSettings.quiz_id == quiz_id).first()
    
    if not qs:
        qs = QuizSettings(
            quiz_id=quiz_id,
            week=payload.week,
            min_difficulty=payload.min_difficulty,
            max_difficulty=payload.max_difficulty,
            min_bloom_level=payload.min_bloom_level,
            max_bloom_level=payload.max_bloom_level,
            max_questions=payload.max_questions,
            allowed_retries=payload.allowed_retries,
            feedback_style=payload.feedback_style,
            include_spaced=payload.include_spaced,
        )
        db.add(qs)
    else:
        qs.week = payload.week
        qs.min_difficulty = payload.min_difficulty
        qs.max_difficulty = payload.max_difficulty
        qs.min_bloom_level = payload.min_bloom_level
        qs.max_bloom_level = payload.max_bloom_level
        qs.max_questions = payload.max_questions
        qs.allowed_retries = payload.allowed_retries
        qs.feedback_style = payload.feedback_style
        qs.include_spaced = payload.include_spaced
        
    db.commit()
    db.refresh(qs)
    
    # ✅ Return properly formatted dict
    return {
        "id": qs.id,
        "quiz_id": qs.quiz_id,
        "quizid": qs.quiz_id,
        "week": qs.week,
        "mindifficulty": qs.min_difficulty,
        "maxdifficulty": qs.max_difficulty,
        "min_bloom_level": qs.min_bloom_level,
        "max_bloom_level": qs.max_bloom_level,
        "maxquestions": qs.max_questions,
        "allowedretries": qs.allowed_retries,
        "feedbackstyle": qs.feedback_style,
        "includespaced": qs.include_spaced
    }




@app.get("/api/quiz/questions/{quiz_id}")
def get_quiz_questions(quiz_id: str, db: Session = Depends(get_db)):
    try:
        mcqs = db.query(MCQ).filter(MCQ.quiz_id == quiz_id).all()
    except AttributeError:
        mcqs = db.query(MCQ).filter(MCQ.quizid == quiz_id).all()
        
    out = []
    for m in mcqs:
        out.append({
            "id": m.id,
            "question": m.question,
            "options": m.options,
            "answer": m.answer,
            "concept_id": m.concept_id,
            "difficulty": m.difficulty or "Medium",
            "bloom_level": m.bloom_level or "Remember"  
        })
        
    return {"quiz_id": quiz_id, "mcqs": out}


@app.post("/api/quiz/generate-title")
async def generate_quiz_title(payload: TitleGenRequest):
    """
    Generates a short, academic quiz title based on a list of concept names.
    """
    if not payload.concepts:
        return {"title": "New Quiz"}

    # Join concepts for the prompt
    topics_str = ", ".join(payload.concepts)
    
    prompt = (
        f"Generate a concise, academic quiz title (3-8 words) for a quiz covering these topics: {topics_str}. "
        "Return ONLY the title text. Do not use quotes. Do not say 'Here is the title'."
        "Examples: 'Introduction to Data Structures', 'Advanced Recursion Patterns', 'Memory Management Basics'."
    )

    try:
        model = genai.GenerativeModel('models/gemini-2.5-flash')
        response = model.generate_content(prompt)
        title = response.text.strip().replace('"', '').replace("'", "")
        return {"title": title}
    except Exception as e:
        print(f"Title Gen Error: {e}")
        # Fallback if LLM fails
        return {"title": f"Quiz: {payload.concepts[0]} & more"}

# --- A. Instructor preview: KG-based sample MCQs (no DB writes) ---
@app.get("/api/quiz/preview")
async def preview_quiz(
    quiz_id: str,
    n: int = 5,
    db: Session = Depends(get_db),
):
    quiz = db.query(Quiz).filter(Quiz.id == quiz_id).first()
    if not quiz:
        raise HTTPException(status_code=404, detail="Quiz not found")

    # Load KG once
    kg_row = db.execute(
        sql_text(
            "SELECT node_data, edge_data FROM knowledge_graph "
            "WHERE course_id=:c AND week=:w"
        ),
        {"c": quiz.course_id, "w": quiz.week},
    ).fetchone()
    kg_nodes = json.loads(kg_row[0]) if (kg_row and kg_row[0]) else []
    kg_edges = json.loads(kg_row[1]) if (kg_row and kg_row[1]) else []

    previews = []

    for cid in quiz.concept_ids:
        selected_node = next((n for n in kg_nodes if n["id"] == cid), None)
        if not selected_node:
            continue

        payload = {
            "course_id": quiz.course_id,
            "week": quiz.week,
            "concept_id": cid,
        }
        resp = await generate_mcqs_kg(payload, db)
        mcqs = resp.get("mcqs", [])[: max(1, n // max(1, len(quiz.concept_ids)))]

        previews.append(
            {
                "concept_id": cid,
                "concept_label": selected_node.get("label", cid),
                "mcqs": mcqs,
            }
        )

    return {"quiz_id": quiz.id, "previews": previews}


@app.post("/api/quiz/lock-preview")
async def lock_preview(payload: LockPreviewPayload, db: Session = Depends(get_db)):
    quiz = db.query(Quiz).filter(Quiz.id == payload.quiz_id).first()
    if not quiz:
        raise HTTPException(status_code=404, detail="Quiz not found")

    created = 0
    for item in payload.items:
        cid = item.get("concept_id")
        for mcq in item.get("mcqs", []):
            if not (
                isinstance(mcq, dict)
                and "question" in mcq
                and "options" in mcq
                and "answer" in mcq
            ):
                continue
            db.add(
                MCQ(
                    id=str(uuid.uuid4()),
                    quiz_id=quiz.id,
                    concept_id=cid,
                    question=mcq["question"],
                    options=mcq["options"],
                    answer=mcq["answer"],
                )
            )
            created += 1

    db.commit()
    return {"quiz_id": quiz.id, "locked_mcqs": created}

@app.post("/api/quiz/generate-mcqs/{quiz_id}")
async def generate_quiz_mcqs(quiz_id: str, payload: GenerateQuizMCQsRequest, db: Session = Depends(get_db)):
    print(f"[DEBUG] Starting generation for Quiz: {quiz_id}")
    
    quiz = db.query(Quiz).filter(Quiz.id == quiz_id).first()
    if not quiz or not quiz.concept_ids:
        return {"quiz_id": quiz_id, "generated_mcqs": 0, "message": "No concepts"}

    count = 0
    target_concepts = quiz.concept_ids[:5]

    for cid in target_concepts:
        kg_payload = {
            "course_id": quiz.course_id,  # Assuming course_id (check model!)
            "week": quiz.week,
            "concept_id": cid,
            "quiz_id": quiz.id
        }
        
        # Handle course_id vs courseid mismatch in payload preparation if needed
        # (The generate_mcqs_kg function we fixed earlier handles the reading part)

        try:
            gen_resp = await generate_mcqs_kg(kg_payload, db)
            new_mcqs = gen_resp.get("mcqs", [])
            print(f"[DEBUG] Generated {len(new_mcqs)} MCQs for concept {cid}")
            
            for m in new_mcqs:
                # Create MCQ object
                new_mcq = MCQ(
                    id=str(uuid.uuid4()),
                    concept_id=cid,
                    question=m.get("question"),
                    options=m.get("options"),
                    answer=m.get("answer"),
                    difficulty=m.get("difficulty", "Medium"), 
                    bloom_level=m.get("bloom_level", "Remember"), 
                )
                
                # Assign FK safely
                if hasattr(MCQ, 'quiz_id'):
                    new_mcq.quiz_id = quiz.id
                else:
                    new_mcq.quizid = quiz.id
                    
                db.add(new_mcq)
                count += 1
        except Exception as e:
            print(f"[ERROR] Failed to save MCQs for {cid}: {e}")
            continue

    db.commit()
    print(f"🎉 Quiz {quiz.id} + {count} MCQs SAVED")
    
    return {"quiz_id": quiz.id, "generated_mcqs": count}



@app.delete("/api/mcq/{mcq_id}")
async def delete_mcq(mcq_id: str, db: Session = Depends(get_db)):
    mcq = db.query(MCQ).filter(MCQ.id == mcq_id).first()
    if not mcq:
        raise HTTPException(status_code=404, detail="MCQ not found")
    db.delete(mcq)
    db.commit()
    return {"status": "deleted", "id": mcq_id}

@app.post("/api/mcq/regenerate/{mcq_id}")
async def regenerate_single_mcq(mcq_id: str, db: Session = Depends(get_db)):
    old_mcq = db.query(MCQ).filter(MCQ.id == mcq_id).first()
    if not old_mcq:
        raise HTTPException(status_code=404, detail="MCQ not found")
    
    quiz = db.query(Quiz).filter(Quiz.id == old_mcq.quiz_id).first()
    
    kg_payload = {
        "course_id": quiz.course_id,
        "week": quiz.week,
        "concept_id": old_mcq.concept_id,
        "quiz_id": quiz.id
    }
    
    try:
        gen_resp = await generate_mcqs_kg(kg_payload, db)
        new_mcqs_data = gen_resp.get("mcqs", [])
        
        if not new_mcqs_data:
             raise HTTPException(status_code=500, detail="Failed to generate replacement")

        best_new = new_mcqs_data[0]
        
        old_mcq.question = best_new["question"]
        old_mcq.options = best_new["options"]
        old_mcq.answer = best_new["answer"]
        old_mcq.difficulty = best_new.get("difficulty", "Medium")
        old_mcq.bloom_level = best_new.get("bloom_level", "Remember")
        
        db.commit()
        return {"status": "regenerated", "mcq": {
            "id": old_mcq.id,
            "question": old_mcq.question,
            "options": old_mcq.options,
            "answer": old_mcq.answer
        }}
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/student/quiz/start", response_model=QuizStartResponse)
async def start_student_quiz(
    quiz_id: str = Body(..., embed=True),
    student_id: str = Body(..., embed=True),
    recommended_bloom_level: Optional[str] = Body(None, embed=True),
    db: Session = Depends(get_db)
):
    # 1. Fetch Quiz & Settings
    quiz = db.query(Quiz).filter(Quiz.id == quiz_id).first()
    if not quiz:
        raise HTTPException(status_code=404, detail="Quiz not found")

    settings = db.query(QuizSettings).filter(QuizSettings.quiz_id == quiz_id).first()

    # Defaults
    max_q = settings.max_questions if settings and settings.max_questions else 10
    allowed_retries = settings.allowed_retries if settings and settings.allowed_retries else 3
    feedback_style = settings.feedback_style if settings and settings.feedback_style else "Immediate"

    # --- Difficulty handling ---
    DIFF_LEVELS = {"easy": 1, "medium": 2, "hard": 3}

    raw_min = settings.min_difficulty if settings and settings.min_difficulty else "Easy"
    raw_max = settings.max_difficulty if settings and settings.max_difficulty else "Hard"

    min_diff_val = DIFF_LEVELS.get(raw_min.lower(), 1)
    max_diff_val = DIFF_LEVELS.get(raw_max.lower(), 3)

    # --- Bloom Levels ---
    BLOOM_LEVELS = ["Remember", "Understand", "Apply", "Analyze", "Evaluate", "Create"]

    raw_min_bloom = settings.min_bloom_level if settings and settings.min_bloom_level else "Remember"
    raw_max_bloom = settings.max_bloom_level if settings and settings.max_bloom_level else "Create"

    try:
        bloom_start = BLOOM_LEVELS.index(raw_min_bloom)
        bloom_end = BLOOM_LEVELS.index(raw_max_bloom)
        if bloom_start > bloom_end:
            bloom_start, bloom_end = bloom_end, bloom_start
        allowed_blooms = set(BLOOM_LEVELS[bloom_start:bloom_end + 1])
    except ValueError:
        allowed_blooms = set(BLOOM_LEVELS)

    # --- Adaptive Bloom targeting ---
    # If the frontend sent a recommended_bloom_level (from /api/student/quiz/adaptive-bloom),
    # validate it falls within the instructor-set range and use it to bias selection.
    adaptive_target_bloom = None
    if recommended_bloom_level and recommended_bloom_level in BLOOM_LEVELS:
        rec_idx = BLOOM_LEVELS.index(recommended_bloom_level)
        bloom_start_idx = BLOOM_LEVELS.index(raw_min_bloom) if raw_min_bloom in BLOOM_LEVELS else 0
        bloom_end_idx = BLOOM_LEVELS.index(raw_max_bloom) if raw_max_bloom in BLOOM_LEVELS else len(BLOOM_LEVELS) - 1
        # Clamp to instructor-allowed range
        clamped_idx = max(bloom_start_idx, min(rec_idx, bloom_end_idx))
        adaptive_target_bloom = BLOOM_LEVELS[clamped_idx]

    print(f"[STUDENT QUIZ] Settings: max_q={max_q}, diff={raw_min}-{raw_max}, bloom={raw_min_bloom}-{raw_max_bloom}")
    print(f"[STUDENT QUIZ] Allowed blooms: {allowed_blooms}, adaptive_target: {adaptive_target_bloom}")

    def get_filtered_mcqs():
        all_mcqs = db.query(MCQ).filter(MCQ.quiz_id == quiz_id).all()
        valid = []
        
        print(f"[FILTER] Total MCQs in quiz: {len(all_mcqs)}")
        
        for m in all_mcqs:
            # Difficulty filter
            q_diff_str = (m.difficulty or "Medium").strip()
            q_val = DIFF_LEVELS.get(q_diff_str.lower(), 2)
            
            if not (min_diff_val <= q_val <= max_diff_val):
                print(f"[FILTER] Rejected MCQ {m.id[:8]} - difficulty {q_diff_str} outside range")
                continue

            # Bloom filter
            q_bloom = (m.bloom_level or "Remember").strip()
            if q_bloom not in allowed_blooms:
                print(f"[FILTER] Rejected MCQ {m.id[:8]} - bloom {q_bloom} not in {allowed_blooms}")
                continue

            print(f"[FILTER] ✅ Accepted MCQ {m.id[:8]} - {q_diff_str}, {q_bloom}")
            valid.append(m)
        
        print(f"[FILTER] Final count: {len(valid)}/{len(all_mcqs)} MCQs passed filters")
        return valid, all_mcqs

    # 2. Find existing attempts for this student+quiz
    attempts_q = db.query(QuizAttempt).filter(
        QuizAttempt.quiz_id == quiz_id,
        QuizAttempt.student_id == student_id
    )

    active_attempt = attempts_q.filter(QuizAttempt.completed == False).first()
    completed_attempts = attempts_q.filter(QuizAttempt.completed == True).count()

    # 3. Resume existing attempt if present
    if active_attempt:
        valid_mcqs, all_mcqs = get_filtered_mcqs()
        target_pool = valid_mcqs if valid_mcqs else all_mcqs

        if not target_pool:
            return {
                "attempt_id": active_attempt.id,
                "questions": [],
                "settings": {
                    "feedback_style": feedback_style,
                    "max_questions": max_q
                },
                "retries_left": max(0, allowed_retries - completed_attempts)
            }

        if adaptive_target_bloom:
            target_bloom_mcqs = [m for m in target_pool if (m.bloom_level or "Remember") == adaptive_target_bloom]
            other_mcqs = [m for m in target_pool if (m.bloom_level or "Remember") != adaptive_target_bloom]
            random.shuffle(target_bloom_mcqs)
            random.shuffle(other_mcqs)
            target_count = min(len(target_bloom_mcqs), max(1, int(max_q * 0.6)))
            other_count = min(len(other_mcqs), max_q - target_count)
            ordered_pool = target_bloom_mcqs[:target_count] + other_mcqs[:other_count]
            if len(ordered_pool) < max_q:
                remaining = [m for m in target_pool if m not in ordered_pool]
                ordered_pool += remaining[:max_q - len(ordered_pool)]
            selected_mcqs = ordered_pool[:max_q]
        else:
            selected_mcqs = random.sample(target_pool, min(len(target_pool), max_q))

        return {
            "attempt_id": active_attempt.id,
            "questions": [
                {
                    "id": m.id,
                    "question": m.question,
                    "options": m.options,
                    "difficulty": m.difficulty,
                    "bloom_level": m.bloom_level
                } for m in selected_mcqs
            ],
            "settings": {
                "feedback_style": feedback_style,
                "max_questions": max_q
            },
            "retries_left": max(0, allowed_retries - completed_attempts)
        }

    # 4. Enforce retry limit
    if completed_attempts >= allowed_retries:
        raise HTTPException(status_code=403, detail="No retries remaining")

    # 5. Create new attempt with filtering
    valid_mcqs, all_mcqs = get_filtered_mcqs()
    target_pool = valid_mcqs if valid_mcqs else all_mcqs

    if not target_pool:
        return {
            "attempt_id": "error",
            "questions": [],
            "settings": {"feedback_style": feedback_style, "max_questions": max_q},
            "retries_left": max(0, allowed_retries - completed_attempts)
        }

    # ✅ Apply adaptive Bloom weighting BEFORE sampling
    # If adaptive_target_bloom is set, front-load questions at that Bloom level,
    # then fill remaining slots with other allowed levels.
    if adaptive_target_bloom:
        target_bloom_mcqs = [m for m in target_pool if (m.bloom_level or "Remember") == adaptive_target_bloom]
        other_mcqs = [m for m in target_pool if (m.bloom_level or "Remember") != adaptive_target_bloom]
        random.shuffle(target_bloom_mcqs)
        random.shuffle(other_mcqs)
        # Aim: ~60% from target bloom level, rest from others (if available)
        target_count = min(len(target_bloom_mcqs), max(1, int(max_q * 0.6)))
        other_count = min(len(other_mcqs), max_q - target_count)
        ordered_pool = target_bloom_mcqs[:target_count] + other_mcqs[:other_count]
        # If we still don't have enough, pad from whatever is left
        if len(ordered_pool) < max_q:
            remaining = [m for m in target_pool if m not in ordered_pool]
            ordered_pool += remaining[:max_q - len(ordered_pool)]
        selected_mcqs = ordered_pool[:max_q]
        print(f"[ADAPTIVE] Targeting '{adaptive_target_bloom}': {target_count} target + {other_count} other = {len(selected_mcqs)} questions")
    else:
        selected_mcqs = random.sample(target_pool, min(len(target_pool), max_q))
    
    print(f"[STUDENT QUIZ] Serving {len(selected_mcqs)} questions to student")

    attempt_id = str(uuid.uuid4())
    new_attempt = QuizAttempt(
        id=attempt_id,
        quiz_id=quiz_id,
        student_id=student_id,
        score=0,
        total_questions=0,
        completed=False
    )
    db.add(new_attempt)
    db.commit()

    return {
        "attempt_id": attempt_id,
        "questions": [
            {
                "id": m.id,
                "question": m.question,
                "options": m.options,
                "difficulty": m.difficulty,
                "bloom_level": m.bloom_level
            } for m in selected_mcqs
        ],
        "settings": {
            "feedback_style": feedback_style,
            "max_questions": max_q
        },
        "retries_left": max(0, allowed_retries - completed_attempts)
    }



@app.put("/api/mcq/{mcq_id}")
async def update_mcq(mcq_id: str, payload: MCQUpdate, db: Session = Depends(get_db)):
    """
    Update an existing MCQ.
    """
    mcq = db.query(MCQ).filter(MCQ.id == mcq_id).first()
    if not mcq:
        raise HTTPException(status_code=404, detail="MCQ not found")
    
    mcq.question = payload.question
    mcq.options = payload.options
    mcq.answer = payload.answer
    mcq.bloom_level = payload.bloom_level 
    if payload.difficulty:
        mcq.difficulty = payload.difficulty
        
    db.commit()
    db.refresh(mcq)
    
    return {
        "status": "updated",
        "mcq": {
            "id": mcq.id,
            "question": mcq.question,
            "options": mcq.options,
            "answer": mcq.answer,
            "difficulty": mcq.difficulty,
            "bloom_level": mcq.bloom_level
        }
    }


@app.post("/api/student/quiz/submit")
async def submit_quiz_attempt(payload: dict = Body(...), db: Session = Depends(get_db)):
    print("SUBMIT payload:", payload)

    attempt_id = payload.get("attempt_id") or payload.get("attemptId")
    student_id = payload.get("student_id") or payload.get("studentId")
    quiz_id = payload.get("quiz_id") or payload.get("quizId")

    raw_responses = payload.get("responses") or payload.get("answers") or []

    responses = []
    for r in raw_responses:
        responses.append({
            "mcq_id": r.get("mcq_id") or r.get("question_id") or r.get("id"),
            "selected_answer": r.get("selected_answer") or r.get("answer") or r.get("selected")
        })

    if not attempt_id or not student_id or not quiz_id:
        raise HTTPException(status_code=400, detail="Missing required fields")

    attempt = db.query(QuizAttempt).filter(QuizAttempt.id == attempt_id).first()
    if not attempt:
        attempt = QuizAttempt(
            id=attempt_id,
            quiz_id=quiz_id,
            student_id=student_id,
            score=0,
            total_questions=0,
            completed=False
        )
        db.add(attempt)
        db.commit()

    correct_count = 0
    total = len(responses)
    per_question_results = []

    for resp in responses:
        mcq_id = resp["mcq_id"]
        selected = resp["selected_answer"]

        mcq = db.query(MCQ).filter(MCQ.id == mcq_id).first()
        if not mcq:
            continue

        is_correct = (selected == mcq.answer)
        if is_correct:
            correct_count += 1

        # ✅ Per-question result for UI
        per_question_results.append({
            "mcq_id": mcq_id,
            "selected": selected,
            "correct": is_correct,
            "correct_answer": mcq.answer,
            "hint": None,
        })

        # ✅ Save/update MCQResponse row (with question)
        existing_response = db.query(MCQResponse).filter(
            MCQResponse.attempt_id == attempt_id,
            MCQResponse.mcq_id == mcq_id
        ).first()

        if existing_response:
            existing_response.selected_answer = selected
            existing_response.is_correct = is_correct
        else:
            new_response = MCQResponse(
                attempt_id=attempt_id,
                mcq_id=mcq_id,
                question=mcq.question,           # ✅ required field
                selected_answer=selected,
                is_correct=is_correct
            )
            db.add(new_response)

    attempt.score = correct_count
    attempt.total_questions = total
    attempt.completed = True
    db.commit()

    # ✅ Broadcast to instructor dashboard
    await manager.broadcast(quiz_id, {
        "type": "submission",
        "student_id": student_id,
        "attempt_id": attempt_id,
        "score": correct_count,
        "total": total,
        "timestamp": datetime.now().isoformat()
    })

    return {
        "success": True,
        "score": correct_count,
        "total": total,
        "percentage": round((correct_count / total * 100) if total > 0 else 0, 1),
        "results": per_question_results
    }

@app.get("/api/quiz/{quiz_id}/questions")
def get_quiz_questions_unfiltered(quiz_id: str, db: Session = Depends(get_db)):
    """
    ✅ Returns ALL MCQs regardless of settings (for instructor question bank)
    """
    mcqs = db.query(MCQ).filter(MCQ.quiz_id == quiz_id).all()
    print(f"[QUESTION BANK] Returning {len(mcqs)} unfiltered MCQs for quiz {quiz_id}")
    
    return [
        {
            "id": m.id,
            "question": m.question,
            "options": m.options,
            "difficulty": m.difficulty or "Medium",
            "bloom_level": m.bloom_level or "Remember"
        } for m in mcqs
    ]


# ✅ Add endpoint to get live quiz stats for instructor
@app.get("/api/quiz/{quiz_id}/live-stats")
async def get_quiz_live_stats(quiz_id: str, db: Session = Depends(get_db)):
    attempts = db.query(QuizAttempt).filter(
        QuizAttempt.quiz_id == quiz_id,
        QuizAttempt.completed == True
    ).all()
    print("LIVE STATS quiz_id:", quiz_id)
    print("All attempt quiz_ids:", [a.quiz_id for a in db.query(QuizAttempt).all()])

    if not attempts:
        return {
            "total_submissions": 0,
            "average_score": 0,
            "submissions": []
        }

    submissions = []
    total_score = 0
    
    for attempt in attempts:
        student = db.query(User).filter(User.id == attempt.student_id).first()
        student_name = f"{student.first_name} {student.last_name}" if student else "Unknown"
        
        submissions.append({
            "student_id": attempt.student_id,
            "student_name": student_name,
            "score": attempt.score,
            "total": attempt.total_questions,
            "percentage": round((attempt.score / attempt.total_questions * 100) if attempt.total_questions > 0 else 0, 1),
            "timestamp": attempt.created_at.isoformat() if hasattr(attempt, 'created_at') else None
        })
        total_score += (attempt.score / attempt.total_questions * 100) if attempt.total_questions > 0 else 0

    return {
        "total_submissions": len(attempts),
        "average_score": round(total_score / len(attempts), 1) if attempts else 0,
        "submissions": submissions
    }


@app.post("/api/student/quiz/check-answer")
async def check_answer_mcq(
    quiz_id: str = Body(..., embed=True),
    mcq_id: str = Body(..., embed=True),
    selected_opt: str = Body(..., embed=True),
    attempt_count: int = Body(1, embed=True),
    db: Session = Depends(get_db)
):
    mcq = db.query(MCQ).filter(MCQ.id == mcq_id).first()
    if not mcq:
        raise HTTPException(status_code=404, detail="Question not found")

    settings = None
    try:
        settings = db.query(QuizSettings).filter(QuizSettings.quiz_id == quiz_id).first()
    except AttributeError:
        settings = db.query(QuizSettings).filter(QuizSettings.quizid == quiz_id).first()

    def get_setting(obj, name, default):
        if not obj: return default
        val = getattr(obj, name, getattr(obj, name.replace("feedback", "feedback_").replace("allowed", "allowed_"), None))
        return val if val is not None else default

    raw_style = get_setting(settings, 'feedbackstyle', "Immediate")
    style = raw_style.lower().strip()
    max_retries = get_setting(settings, 'allowedretries', 3)

    stored_ans = (mcq.answer or "").strip()
    user_ans = (selected_opt or "").strip()
    
    is_correct = stored_ans.lower() == user_ans.lower()

    if not is_correct and len(stored_ans) == 1 and stored_ans.upper() in ['A', 'B', 'C', 'D']:
        try:
            options_list = mcq.options
            if isinstance(options_list, str):
                options_list = json.loads(options_list)
            
            if isinstance(options_list, list):
                idx = ord(stored_ans.upper()) - 65
                if 0 <= idx < len(options_list):
                    correct_text = str(options_list[idx]).strip()
                    if correct_text.lower() == user_ans.lower():
                        is_correct = True
        except Exception as e:
            print(f"Legacy answer check failed: {e}")

    response = {
        "correct": is_correct,
        "style": raw_style,
        "retries_exhausted": False
    }

    if not is_correct:
        if attempt_count >= max_retries:
            response["retries_exhausted"] = True
            response["correct_answer"] = mcq.answer 
            response["hint"] = "No more retries."
        else:
            if "hint" in style:
                concept_text = mcq.concept_id if mcq.concept_id else "course concepts"
                response["hint"] = f"Review {concept_text}"
            else:
                 response["hint"] = "Incorrect. Try again."

    return response


async def generate_knowledge_graph(course_id: str, week: int, processing_id: str, db: Session):
    """
    Analyzes lecture slides to build a hierarchical concept map.
    Forces 'main_topic' to be the root node.
    """
    try:
        print(f"[KG START] Generating graph for Course: {course_id}, Week: {week}")
        
        # Get upload record for status updates
        upload_record = db.query(LectureProcessing).filter(LectureProcessing.id == processing_id).first()

        # 1. Fetch Text Content
        segments = db.query(Segment).filter(
            Segment.course_id == course_id, 
            Segment.week == week
        ).all()
        
        if not segments:
            print("[KG SKIP] No segments found.")
            return

        full_text = "\n".join([s.content for s in segments if s.content])[:15000]
        file_name = upload_record.file_name if upload_record else "Lecture Topic"

        # 2. Call LLM
        model = genai.GenerativeModel("gemini-2.5-flash")
        
        prompt = f"""
        You are an expert computer science instructor. Analyze the following lecture content from "{file_name}".
        
        STEP 1: Identify the SINGLE main topic (e.g., "Stacks", "Queue").
        STEP 2: Identify sub-concepts and their relationships.

        Return valid JSON:
        {{
            "main_topic": "The single overarching topic string",
            "nodes": [ {{"id": "Concept Name", "label": "Display Label"}} ],
            "edges": [ {{"source": "Parent ID", "target": "Child ID", "relation": "..."}} ]
        }}

        RULES:
        - "main_topic" is the absolute root.
        - All concepts must connect back to "main_topic".
        - No disconnected islands.

        Context:
        {full_text}
        """

        response = model.generate_content(prompt)
        # robust json cleaning
        cleaned_text = response.text.replace("```json", "").replace("```", "").strip()
        data = json.loads(cleaned_text)

        # 3. Process Hierarchy
        main_topic = data.get("main_topic", file_name.replace(".pdf", ""))
        raw_nodes = data.get("nodes", [])
        raw_edges = data.get("edges", [])

        print(f"[KG LOGIC] Main Topic: {main_topic}")

        # Use the compute_levels helper we defined earlier
        final_nodes, final_edges = compute_levels(raw_nodes, raw_edges, explicit_root=main_topic)

        # 4. Save to DB
        existing_kg = db.query(KnowledgeGraphBase).filter(
            KnowledgeGraphBase.course_id == course_id,
            KnowledgeGraphBase.week == week
        ).first()

        node_json = json.dumps(final_nodes)
        edge_json = json.dumps(final_edges)

        if existing_kg:
            existing_kg.node_data = node_json
            existing_kg.edge_data = edge_json
        else:
            new_kg = KnowledgeGraphBase(
                id=str(uuid.uuid4()),
                course_id=course_id,
                week=week,
                node_data=node_json,
                edge_data=edge_json
            )
            db.add(new_kg)

        # 5. Mark Processing Complete
        if upload_record:
            upload_record.status = "done"
            upload_record.progress = 100
        
        db.commit()
        print("[KG SUCCESS] Graph saved.")

    except Exception as e:
        print(f"[KG ERROR] {str(e)}")
        traceback.print_exc()
        if upload_record:
            upload_record.status = "error"
            upload_record.error_message = str(e)
            db.commit()

# ===== STUDENT PERFORMANCE + INSTRUCTOR FEEDBACK ENDPOINTS =====


# ---------------------------------------------------------------------------
# ENDPOINT 1: GET /api/student/performance
# ---------------------------------------------------------------------------

@app.get("/api/student/performance")
def get_student_performance(student_id: str, course_id: str, db: Session = Depends(get_db)):
    # Fetch all completed attempts for this student in the given course
    attempts = (
        db.query(QuizAttempt)
        .join(Quiz, Quiz.id == QuizAttempt.quiz_id)
        .filter(
            QuizAttempt.student_id == student_id,
            QuizAttempt.completed == True,
            Quiz.course_id == course_id,
        )
        .order_by(QuizAttempt.started_at.desc())
        .all()
    )

    total_attempts = len(attempts)
    overall_score_pct = 0.0
    quiz_history = []

    bloom_levels = ["Remember", "Understand", "Apply", "Analyze", "Evaluate", "Create"]
    bloom_breakdown = {lvl: {"correct": 0, "total": 0, "pct": 0.0} for lvl in bloom_levels}
    concept_map = {}  # concept_id -> {"correct": int, "total": int}

    score_pcts = []

    for attempt in attempts:
        quiz = db.query(Quiz).filter(Quiz.id == attempt.quiz_id).first()
        quiz_name = quiz.title if quiz and hasattr(quiz, "title") else str(attempt.quiz_id)

        total = attempt.total_questions if attempt.total_questions else 0
        score = attempt.score if attempt.score is not None else 0
        pct = round((score / total * 100), 2) if total > 0 else 0.0
        score_pcts.append(pct)

        # Derive week number from started_at
        week_num = None
        if attempt.started_at:
            week_num = attempt.started_at.isocalendar()[1]

        quiz_history.append({
            "attempt_id": str(attempt.id),
            "quiz_id": str(attempt.quiz_id),
            "quiz_name": quiz_name,
            "score": score,
            "total": total,
            "score_pct": pct,
            "completed_at": attempt.started_at.isoformat() if attempt.started_at else None,
            "week": week_num,
        })

        # Aggregate bloom and concept data from MCQResponses
        responses = (
            db.query(MCQResponse)
            .filter(MCQResponse.attempt_id == attempt.id)
            .all()
        )

        for resp in responses:
            mcq = db.query(MCQ).filter(MCQ.id == resp.mcq_id).first()
            if not mcq:
                continue

            bloom = mcq.bloom_level
            concept = str(mcq.concept_id) if mcq.concept_id else "unknown"

            if bloom in bloom_breakdown:
                bloom_breakdown[bloom]["total"] += 1
                if resp.is_correct:
                    bloom_breakdown[bloom]["correct"] += 1

            if concept not in concept_map:
                concept_map[concept] = {"correct": 0, "total": 0}
            concept_map[concept]["total"] += 1
            if resp.is_correct:
                concept_map[concept]["correct"] += 1

    # Compute overall score pct
    if score_pcts:
        overall_score_pct = round(sum(score_pcts) / len(score_pcts), 2)

    # Finalise bloom breakdown percentages
    for lvl in bloom_levels:
        t = bloom_breakdown[lvl]["total"]
        c = bloom_breakdown[lvl]["correct"]
        bloom_breakdown[lvl]["pct"] = round((c / t * 100), 2) if t > 0 else 0.0

    # Build concept mastery list
    concept_mastery = []
    for concept_id, data in concept_map.items():
        t = data["total"]
        c = data["correct"]
        pct = round((c / t * 100), 2) if t > 0 else 0.0
        if pct >= 80:
            status = "mastered"
        elif pct >= 50:
            status = "learning"
        else:
            status = "struggling"
        concept_mastery.append({
            "concept_id": concept_id,
            "correct": c,
            "total": t,
            "pct": pct,
            "status": status,
        })

    # Sort weakest first
    concept_mastery.sort(key=lambda x: x["pct"])

    return {
        "total_attempts": total_attempts,
        "overall_score_pct": overall_score_pct,
        "quiz_history": quiz_history,
        "bloom_breakdown": bloom_breakdown,
        "concept_mastery": concept_mastery,
    }


# ---------------------------------------------------------------------------
# ENDPOINT 2: GET /api/student/quiz/adaptive-bloom
# ---------------------------------------------------------------------------

@app.get("/api/student/quiz/adaptive-bloom")
def get_adaptive_bloom(student_id: str, quiz_id: str, db: Session = Depends(get_db)):
    bloom_levels = ["Remember", "Understand", "Apply", "Analyze", "Evaluate", "Create"]

    # Fetch quiz settings
    settings = db.query(QuizSettings).filter(QuizSettings.quiz_id == quiz_id).first()
    min_bloom = settings.min_bloom_level if settings and settings.min_bloom_level else bloom_levels[0]
    max_bloom = settings.max_bloom_level if settings and settings.max_bloom_level else bloom_levels[-1]

    # Clamp allowed range
    min_idx = bloom_levels.index(min_bloom) if min_bloom in bloom_levels else 0
    max_idx = bloom_levels.index(max_bloom) if max_bloom in bloom_levels else len(bloom_levels) - 1
    allowed_levels = bloom_levels[min_idx: max_idx + 1]

    # Fetch all completed attempts for this student + quiz
    attempts = (
        db.query(QuizAttempt)
        .filter(
            QuizAttempt.student_id == student_id,
            QuizAttempt.quiz_id == quiz_id,
            QuizAttempt.completed == True,
        )
        .all()
    )

    # Aggregate per-bloom accuracy
    bloom_performance = {lvl: {"correct": 0, "total": 0, "pct": 0.0} for lvl in bloom_levels}

    for attempt in attempts:
        responses = (
            db.query(MCQResponse)
            .filter(MCQResponse.attempt_id == attempt.id)
            .all()
        )
        for resp in responses:
            mcq = db.query(MCQ).filter(MCQ.id == resp.mcq_id).first()
            if not mcq:
                continue
            bloom = mcq.bloom_level
            if bloom in bloom_performance:
                bloom_performance[bloom]["total"] += 1
                if resp.is_correct:
                    bloom_performance[bloom]["correct"] += 1

    for lvl in bloom_levels:
        t = bloom_performance[lvl]["total"]
        c = bloom_performance[lvl]["correct"]
        bloom_performance[lvl]["pct"] = round((c / t * 100), 2) if t > 0 else 0.0

    # Determine recommended bloom level
    recommended = None
    for lvl in allowed_levels:
        data = bloom_performance[lvl]
        # Not tested yet OR accuracy below threshold
        if data["total"] == 0 or data["pct"] < 75.0:
            recommended = lvl
            break

    # If all levels are >= 75%, return the highest allowed level
    if recommended is None:
        recommended = allowed_levels[-1]

    # Build message
    recommended_idx = bloom_levels.index(recommended)
    if recommended_idx > 0:
        prev_level = bloom_levels[recommended_idx - 1]
        message = f"Targeting {recommended} — keep practicing {prev_level}"
    else:
        message = f"Targeting {recommended} — great starting point!"

    return {
        "recommended_bloom_level": recommended,
        "bloom_performance": bloom_performance,
        "message": message,
    }



# ===== INSTRUCTOR ENROLLMENT MANAGEMENT + ANONYMOUS FEEDBACK =====

# ── Endpoint 1: GET /api/users/students ──────────────────────────────────────

@app.get("/api/users/students")
def get_all_students(
    search: Optional[str] = None,
    db: Session = Depends(get_db),
):
    query = db.query(User).filter(User.role == "student")

    if search:
        s = search.lower()
        query = query.filter(
            or_(
                func.lower(User.email).contains(s),
                func.lower(User.id).startswith(s),
            )
        )

    students = query.limit(50).all()

    return [
        {
            "id": u.id,
            "email": u.email,
            "truncated_id": u.id[:8],
        }
        for u in students
    ]


# ── Endpoint 2: GET /api/course/{course_id}/students ─────────────────────────

@app.get("/api/course/{course_id}/students")
def get_course_students(
    course_id: str,
    db: Session = Depends(get_db),
):
    rows = (
        db.query(Enrollment, User)
        .join(User, Enrollment.student_id == User.id)
        .filter(Enrollment.course_id == course_id)
        .all()
    )

    return [
        {
            "student_id": enrollment.student_id,
            "email": user.email,
            "enrolled_at": enrollment.enrolled_at.isoformat() if enrollment.enrolled_at else None,
            "enrollment_id": enrollment.id,
        }
        for enrollment, user in rows
    ]


# ── Endpoint 3: POST /api/instructor/enroll ──────────────────────────────────

@app.post("/api/instructor/enroll")
def enroll_student(
    body: dict = Body(...),
    db: Session = Depends(get_db),
):
    course_id = body.get("course_id")
    instructor_id = body.get("instructor_id")
    identifier = body.get("identifier")

    # 1. Verify instructor owns course
    course = db.query(Course).filter(Course.id == course_id).first()
    if not course or course.instructor_id != instructor_id:
        raise HTTPException(status_code=403, detail="Not authorized for this course")

    # 2. Look up student by UUID first, then by email
    user = db.query(User).filter(User.id == identifier).first()
    if not user:
        user = db.query(User).filter(func.lower(User.email) == identifier.lower()).first()

    # 3. No user found
    if not user:
        raise HTTPException(status_code=404, detail="No student found with that ID or email")

    # 4. Must be a student account
    if user.role != "student":
        raise HTTPException(status_code=400, detail="That account is not a student")

    # 5. Already enrolled?
    existing = (
        db.query(Enrollment)
        .filter(Enrollment.course_id == course_id, Enrollment.student_id == user.id)
        .first()
    )
    if existing:
        raise HTTPException(status_code=409, detail="Student already enrolled")

    # 6. Create enrollment
    enrollment = Enrollment(
        id=str(uuid.uuid4()),
        course_id=course_id,
        student_id=user.id,
    )
    db.add(enrollment)
    db.commit()
    db.refresh(enrollment)

    return {
        "status": "enrolled",
        "student_id": user.id,
        "email": user.email,
        "enrollment_id": enrollment.id,
    }


# ── Endpoint 4: DELETE /api/instructor/unenroll ──────────────────────────────

@app.delete("/api/instructor/unenroll")
def unenroll_student(
    body: dict = Body(...),
    db: Session = Depends(get_db),
):
    course_id = body.get("course_id")
    instructor_id = body.get("instructor_id")
    student_id = body.get("student_id")

    # 1. Verify instructor owns course
    course = db.query(Course).filter(Course.id == course_id).first()
    if not course or course.instructor_id != instructor_id:
        raise HTTPException(status_code=403, detail="Not authorized for this course")

    # 2. Find and delete enrollment
    enrollment = (
        db.query(Enrollment)
        .filter(Enrollment.course_id == course_id, Enrollment.student_id == student_id)
        .first()
    )
    if not enrollment:
        raise HTTPException(status_code=404, detail="Enrollment not found")

    db.delete(enrollment)
    db.commit()

    return {"status": "removed", "student_id": student_id}


# ── Endpoint 5: GET /api/instructor/quiz/feedback (aggregate/anonymous) ──────

@app.get("/api/instructor/quiz/feedback")
def get_quiz_feedback_aggregate(
    quiz_id: str,
    db: Session = Depends(get_db),
):
    quiz = db.query(Quiz).filter(Quiz.id == quiz_id).first()
    if not quiz:
        raise HTTPException(status_code=404, detail="Quiz not found")

    all_attempts = db.query(QuizAttempt).filter(QuizAttempt.quiz_id == quiz_id).all()
    total_attempts = len(all_attempts)
    completed_attempts = [a for a in all_attempts if a.completed]
    completed_count = len(completed_attempts)

    completion_rate = (completed_count / total_attempts * 100) if total_attempts > 0 else 0.0

    # Average score across completed attempts only
    scores = [a.score_pct for a in completed_attempts if a.score_pct is not None]
    avg_score_pct = (sum(scores) / len(scores)) if scores else 0.0

    # Score distribution buckets
    distribution = {"0-20": 0, "21-40": 0, "41-60": 0, "61-80": 0, "81-100": 0}
    for s in scores:
        if s <= 20:
            distribution["0-20"] += 1
        elif s <= 40:
            distribution["21-40"] += 1
        elif s <= 60:
            distribution["41-60"] += 1
        elif s <= 80:
            distribution["61-80"] += 1
        else:
            distribution["81-100"] += 1

    # Question-level breakdown
    mcqs = db.query(MCQ).filter(MCQ.quiz_id == quiz_id).all()
    question_breakdown = []

    bloom_accumulator = {}  # bloom_level -> {"questions": int, "total_accuracy": float}

    for mcq in mcqs:
        responses = (
            db.query(MCQResponse)
            .join(QuizAttempt, MCQResponse.attempt_id == QuizAttempt.id)
            .filter(MCQResponse.mcq_id == mcq.id, QuizAttempt.completed == True)
            .all()
        )

        total_answers = len(responses)
        correct_count = sum(1 for r in responses if r.is_correct)
        accuracy_pct = (correct_count / total_answers * 100) if total_answers > 0 else 0.0

        # Most common wrong answer
        wrong_answers = [r.selected_answer for r in responses if not r.is_correct and r.selected_answer]
        most_common_wrong = None
        if wrong_answers:
            counter = Counter(wrong_answers)
            most_common_wrong = counter.most_common(1)[0][0]

        bloom_level = mcq.bloom_level or "Unknown"

        question_breakdown.append({
            "mcq_id": mcq.id,
            "question": mcq.question,
            "bloom_level": bloom_level,
            "difficulty": mcq.difficulty,
            "concept_id": mcq.concept_id,
            "total_answers": total_answers,
            "correct_count": correct_count,
            "accuracy_pct": round(accuracy_pct, 2),
            "most_common_wrong_answer": most_common_wrong,
        })

        # Accumulate bloom stats
        if bloom_level not in bloom_accumulator:
            bloom_accumulator[bloom_level] = {"questions": 0, "total_accuracy": 0.0}
        bloom_accumulator[bloom_level]["questions"] += 1
        bloom_accumulator[bloom_level]["total_accuracy"] += accuracy_pct

    # Sort by accuracy ascending (hardest first)
    question_breakdown.sort(key=lambda x: x["accuracy_pct"])

    # Build bloom_summary for all six levels
    bloom_levels = ["Remember", "Understand", "Apply", "Analyze", "Evaluate", "Create"]
    bloom_summary = {}
    for level in bloom_levels:
        if level in bloom_accumulator:
            acc = bloom_accumulator[level]
            n = acc["questions"]
            bloom_summary[level] = {
                "questions": n,
                "avg_accuracy": round(acc["total_accuracy"] / n, 2) if n > 0 else 0.0,
            }
        else:
            bloom_summary[level] = {"questions": 0, "avg_accuracy": 0.0}

    return {
        "quiz_id": quiz_id,
        "quiz_name": quiz.name,
        "total_attempts": total_attempts,
        "completed_attempts": completed_count,
        "completion_rate": round(completion_rate, 2),
        "avg_score_pct": round(avg_score_pct, 2),
        "score_distribution": distribution,
        "question_breakdown": question_breakdown,
        "bloom_summary": bloom_summary,
    }
